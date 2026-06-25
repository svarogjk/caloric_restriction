"""Generate a 45-slide PowerPoint presentation for PyData Yerevan (80 min, SE audience).

Section 1 — Hook & Problem (slides 1–2)
  1.  Title
  2.  The Problem: 200K+ Datasets Nobody Can Use

Section 2 — Domain Primer (slides 3–6)
  3.  What Is Survival Analysis?
  4.  Kaplan-Meier Curves
  5.  Hazard Ratio: One Number to Rule Them All
  6.  Meta-Analysis: Why N Datasets > 1 Dataset

Section 3 — What We Built (slide 7)
  7.  What We Built: NL Query → Results in 5 Min

Section 4 — Backend Deep Dive (slides 8–15)
  8.  Backend Architecture  (visual: routes → services schema)
  9.  FastAPI: 4 Problems Solved for This App
  10. Pydantic: One Model, Four Jobs
  11. asyncio: Production Patterns Beyond gather()
  12. pydantic-ai: Production-Grade AI Agents
  13. lifelines: Survival Analysis in Python
  14. SQLAlchemy Async: One Codebase, Two Databases
  15. Mistral Embeddings + uv: RAG Without Infrastructure

Section 5 — Frontend Deep Dive (slides 16–20)
  16. Frontend Architecture  (visual: component tree + Redux + API)
  17. React 18: Component-Based UI Framework
  18. Redux Toolkit: Predictable Global State
  19. Server-Sent Events: Real-Time Analysis Progress
  20. Recharts: Declarative Data Visualisation

Section 6 — Pipeline & Services (slides 21–29)
  21. The 6-Step Workflow We Automate
  22. GEOClient: Accessing 200K+ Datasets from NCBI GEO
  23. GEOLoaderService: Parsing Expression Matrices
  24. GEOSurvivalWorkflowOrchestrator: Pipeline Coordinator
  25. GEORankingService: Selecting the Best Datasets
  26. Gene Mapping: The Platform Babel Problem
  27. Gene Mapping: Caching Makes It Possible
  28. Statistical Engine: lifelines (overview)
  29. SurvivalAnalysisService: Statistical Core

Section 7 — AI Chat (slides 30–33)
  30. AI Chat: pydantic-ai + RAG + Domain Score (overview)
  31. AI Chat: 5 Tools Grounding Every Response in Real Data
  32. RAG: Retrieval-Augmented Generation with numpy
  33. Domain Score: Measuring AI Response Quality

Section 8 — Cross-Cutting Concerns (slides 34–39)
  34. End-to-End SSE Pipeline: From Query to Live Results
  35. JWT Authentication: Stateless Security
  36. Full Data Pipeline: NL Query → Publication Export
  37. Multi-Layer Caching: Architecture Overview
  38. asyncio.gather: The 10× Speedup in Practice
  39. Error Handling: Resilient by Design

Section 9 — Engineering Decisions (slides 40–41)
  40. SQLite vs PostgreSQL: Cost vs Capability
  41. Deployment Pipeline

Section 10 — Demo (slides 42–44)
  42. Demo: Natural Language Query
  43. Demo: Volcano Plot & Gene Results
  44. Demo: Kaplan-Meier Survival Curves

Section 11 — Lessons Learned (slide 45)
  45. What We Learned
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x19, 0x23)
CARD      = RGBColor(0x1A, 0x27, 0x36)
CARD_ALT  = RGBColor(0x15, 0x20, 0x30)
WHITE     = RGBColor(0xF0, 0xF0, 0xF0)
MUTED     = RGBColor(0xA0, 0xAE, 0xC0)
LIGHT     = RGBColor(0xC8, 0xD6, 0xE5)
TEAL      = RGBColor(0x2D, 0xD4, 0xA8)
BLUE      = RGBColor(0x60, 0xA5, 0xFA)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
AMBER     = RGBColor(0xFB, 0xBF, 0x24)
CORAL     = RGBColor(0xF9, 0x70, 0x66)
TEAL_DIM  = RGBColor(0x1A, 0x5E, 0x4E)
CORAL_DIM = RGBColor(0x6B, 0x2E, 0x2A)
TEAL_DARK = RGBColor(0x1A, 0x40, 0x35)

# ── Primitive helpers ──────────────────────────────────────────────────────────

def _set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _rect(slide, left, top, width, height, fill_color=CARD, name: str = "") -> object:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _oval(slide, left, top, width, height, fill_color=TEAL, name: str = "") -> object:
    shape = slide.shapes.add_shape(9, Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _tb(slide, left: float, top: float, width: float, height: float, name: str = "") -> object:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        txBox.name = name
    txBox.text_frame.word_wrap = True
    return txBox


def _para(
    tf,
    text: str,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    new: bool = False,
    space_after: int = 0,
    mono: bool = False,
) -> None:
    p = tf.add_paragraph() if new else tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if mono:
        p.font.name = "Courier New"


def _title_block(slide, title: str, accent: RGBColor = TEAL) -> None:
    tb = _tb(slide, 0.5, 0.30, 12.333, 0.75, "Title")
    _para(tb.text_frame, title, 28, WHITE, bold=True)
    _rect(slide, 0.5, 0.95, 12.333, 0.04, fill_color=accent, name="TitleLine")


# ── Slide 1: Title ─────────────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, title: str, sub1: str, sub2: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _oval(slide, 10.0, -0.5, 4.0, 4.0, TEAL,   "DecorCircle1")
    _oval(slide, -1.0,  5.0, 3.5, 3.5, BLUE,   "DecorCircle2")
    _rect(slide, 5.25, 2.40, 2.833, 0.06, TEAL, "AccentLine")
    tb = _tb(slide, 1.5, 2.60, 10.333, 1.1, "MainTitle")
    _para(tb.text_frame, title, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb2 = _tb(slide, 2.0, 3.89, 9.333, 1.2, "Subtitle")
    _para(tb2.text_frame, sub1,  20, TEAL,  align=PP_ALIGN.CENTER)
    _para(tb2.text_frame, sub2,  16, MUTED, align=PP_ALIGN.CENTER, new=True)


# ── Slide 2: Combined intro (Problem) ─────────────────────────────────────────

def add_combined_intro_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "The Problem: 200K+ Datasets Nobody Can Use", CORAL)

    tb_ph = _tb(slide, 0.5, 1.10, 5.80, 0.35, "ProblemHeader")
    _para(tb_ph.text_frame, "THE SCALE", 12, CORAL, bold=True)

    stat_data = [
        ("$2.6B",  "Per Approved Drug",  AMBER),
        ("200K+",  "Open Datasets in GEO", TEAL),
        ("6 Steps", "Manual Expert Work",  CORAL),
    ]
    stat_w, stat_h = 1.80, 0.85
    stat_xs = [0.5, 2.40, 4.30]
    stat_y = 1.50
    for idx, (stat, label, color) in enumerate(stat_data):
        x = stat_xs[idx]
        _rect(slide, x, stat_y, stat_w, stat_h, CARD, f"StatBadge{idx}")
        _rect(slide, x, stat_y, 0.06, stat_h, color, f"StatAccent{idx}")
        tb_s = _tb(slide, x + 0.10, stat_y + 0.05, stat_w - 0.15, 0.42)
        _para(tb_s.text_frame, stat, 24, color, bold=True, align=PP_ALIGN.CENTER)
        tb_l = _tb(slide, x + 0.10, stat_y + 0.50, stat_w - 0.15, 0.30)
        _para(tb_l.text_frame, label, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb_gh = _tb(slide, 0.5, 2.45, 5.80, 0.35, "GapsHeader")
    _para(tb_gh.text_frame, "WHY IT'S HARD", 12, AMBER, bold=True)

    gap_data = [
        ("Data Silos",        "Internal pharma data locked by IP, GxP compliance",  CORAL),
        ("Platform Babel",    "100+ microarray platforms with proprietary probe IDs", TEAL),
        ("Single-Study Bias", "Tools limited to TCGA — no cross-cohort validation",  CORAL),
        ("No NL Interface",   "Requires R/Python expertise to access GEO",           PURPLE),
        ("Manual Curation",   "Hours of expert work per dataset to extract survival", AMBER),
        ("Tool Limits",       "KMplot / GEPIA2 locked to TCGA, cannot search GEO",   BLUE),
    ]
    gap_xs = [0.5, 3.20]
    gap_ys = [2.85, 3.65, 4.45]
    gw, gh = 2.60, 0.70
    for idx, (heading, desc, color) in enumerate(gap_data):
        col = idx % 2
        row = idx // 2
        x = gap_xs[col]
        y = gap_ys[row]
        _rect(slide, x, y, gw, gh, CARD, f"GapBadge{idx}")
        _rect(slide, x, y, 0.06, gh, color, f"GapAccent{idx}")
        tb_h = _tb(slide, x + 0.10, y + 0.05, gw - 0.15, 0.28)
        _para(tb_h.text_frame, heading, 11, color, bold=True)
        tb_d = _tb(slide, x + 0.10, y + 0.34, gw - 0.15, 0.32)
        _para(tb_d.text_frame, desc, 10, MUTED)

    tb_src = _tb(slide, 0.5, 5.22, 5.80, 0.35, "Sources")
    _para(tb_src.text_frame,
          "Sources: DiMasi et al. JHPM 2016  ·  FDA 21 CFR Part 11  ·  NCBI GEO 2024",
          8, MUTED)

    _rect(slide, 6.35, 1.05, 0.04, 5.20, TEAL, "Divider")

    tb_ah = _tb(slide, 6.55, 1.10, 6.30, 0.35, "ArchHeader")
    _para(tb_ah.text_frame, "FULL SYSTEM ARCHITECTURE", 12, TEAL, bold=True)

    arch_layers = [
        ("Client",       ["React 18 + Redux", "Recharts / SSE"],                   BLUE),
        ("API",          ["FastAPI Router",    "JWT + SSE stream"],                 TEAL),
        ("Orchestrator", ["GEOSurvivalWorkflow", "asyncio.gather"],                 PURPLE),
        ("Services",     ["GEOClient", "Survival / Ranking", "GeneMapping"],        AMBER),
        ("Storage",      ["SQLite / PostgreSQL", "Parquet cache"],                  MUTED),
    ]
    row_h   = 0.82
    arrow_h = 0.13
    start_y = 1.55
    lbl_x   = 6.55
    lbl_w   = 1.05
    box_x   = 7.70
    box_w   = 4.90
    gap_box = 0.10

    for i, (label, boxes, color) in enumerate(arch_layers):
        y = start_y + i * (row_h + arrow_h)
        tb_lbl = _tb(slide, lbl_x, y + 0.22, lbl_w, 0.35, f"ArchLabel{i}")
        _para(tb_lbl.text_frame, label, 9, MUTED, align=PP_ALIGN.RIGHT)
        n  = len(boxes)
        bw = (box_w - (n - 1) * gap_box) / n
        for j, text in enumerate(boxes):
            bx = box_x + j * (bw + gap_box)
            _rect(slide, bx, y, bw, row_h, CARD, f"ArchBox{i}_{j}")
            _rect(slide, bx, y, bw, 0.06, color, f"ArchStripe{i}_{j}")
            tb_b = _tb(slide, bx, y + 0.16, bw, row_h - 0.22, f"ArchText{i}_{j}")
            _para(tb_b.text_frame, text, 11, LIGHT, align=PP_ALIGN.CENTER)
        if i < len(arch_layers) - 1:
            ax = box_x + box_w / 2 - 0.06
            _rect(slide, ax, y + row_h, 0.12, arrow_h, TEAL, f"ArchArrow{i}")


# ── Slide 3: What Is Survival Analysis? ───────────────────────────────────────

def add_survival_primer_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "What Is Survival Analysis?", TEAL)

    # Patient data: (end_time_months, is_event, is_high_expression)
    patients = [
        (8,  True,  True),
        (15, True,  True),
        (22, True,  True),
        (33, True,  True),
        (42, False, True),
        (48, True,  False),
        (52, False, False),
        (55, False, False),
        (57, True,  False),
        (60, False, False),
    ]

    cl   = 1.60
    ct   = 1.25
    cw   = 6.50
    ch   = 4.40
    xmax = 60
    xs   = cw / xmax
    row_h = ch / len(patients)
    bar_h = 0.18

    # X-axis baseline
    _rect(slide, cl, ct + ch, cw, 0.03, MUTED, "XAxis")

    for i, (end_t, is_event, is_high) in enumerate(patients):
        color    = CORAL if is_high else TEAL
        row_top  = ct + i * row_h
        bar_top  = row_top + (row_h - bar_h) / 2.0
        bar_w    = end_t * xs

        _rect(slide, cl, bar_top, bar_w, bar_h, color, f"PBar{i}")

        mx = cl + bar_w
        if is_event:
            _rect(slide, mx - 0.07, bar_top - 0.07, 0.14, bar_h + 0.14, color, f"PEvent{i}")
        else:
            _rect(slide, mx - 0.03, bar_top - 0.12, 0.06, bar_h + 0.24, WHITE, f"PCensor{i}")

    # X-axis ticks and labels
    for t in [0, 12, 24, 36, 48, 60]:
        tx = cl + t * xs
        _rect(slide, tx - 0.02, ct + ch, 0.04, 0.10, MUTED, f"XTick{t}")
        tb_t = _tb(slide, tx - 0.28, ct + ch + 0.11, 0.56, 0.30)
        _para(tb_t.text_frame, str(t), 10, MUTED, align=PP_ALIGN.CENTER)

    tb_xt = _tb(slide, cl, ct + ch + 0.45, cw, 0.28)
    _para(tb_xt.text_frame, "Time (months)", 12, MUTED, align=PP_ALIGN.CENTER)

    # Panel header
    tb_ph = _tb(slide, cl, 1.00, 4.0, 0.25)
    _para(tb_ph.text_frame, "10 PATIENTS — 2 EXPRESSION GROUPS", 11, MUTED, bold=True)

    # Legend
    _rect(slide, cl, 6.30, 0.30, 0.18, CORAL, "LegHighRect")
    tb_lh = _tb(slide, cl + 0.36, 6.26, 2.5, 0.25)
    _para(tb_lh.text_frame, "High expression  ●=event", 10, CORAL)
    _rect(slide, cl + 3.2, 6.30, 0.30, 0.18, TEAL, "LegLowRect")
    tb_ll = _tb(slide, cl + 3.56, 6.26, 2.5, 0.25)
    _para(tb_ll.text_frame, "Low expression  |=censored", 10, TEAL)

    # ── Right panel: Concepts ─────────────────────────────────────────────────
    rx = 8.50
    rw = 12.633 - rx

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "TWO THINGS WE MEASURE", 12, MUTED, bold=True)

    # EVENT card
    _rect(slide, rx, 1.50, rw, 1.55, CARD, "EventCard")
    _rect(slide, rx, 1.50, 0.08, 1.55, CORAL, "EventAccent")
    tb_eh = _tb(slide, rx + 0.18, 1.60, rw - 0.28, 0.38)
    _para(tb_eh.text_frame, "EVENT  ●", 16, CORAL, bold=True)
    tb_ed = _tb(slide, rx + 0.18, 2.05, rw - 0.28, 0.88)
    _para(tb_ed.text_frame,
          "Death, relapse, or disease progression.\nThe outcome we are studying.", 13, LIGHT)

    # CENSORED card
    _rect(slide, rx, 3.20, rw, 1.55, CARD, "CensorCard")
    _rect(slide, rx, 3.20, 0.08, 1.55, TEAL, "CensorAccent")
    tb_ch = _tb(slide, rx + 0.18, 3.30, rw - 0.28, 0.38)
    _para(tb_ch.text_frame, "CENSORED  |", 16, TEAL, bold=True)
    tb_cd = _tb(slide, rx + 0.18, 3.75, rw - 0.28, 0.90)
    _para(tb_cd.text_frame,
          "Last known contact, lost to follow-up, study ended. Patient was alive at last observation.", 13, LIGHT)

    # Key insight
    _rect(slide, rx, 4.90, rw, 1.65, CARD, "InsightCard")
    _rect(slide, rx, 4.90, rw, 0.05, TEAL, "InsightTop")
    tb_ik = _tb(slide, rx + 0.15, 5.03, rw - 0.25, 0.32)
    _para(tb_ik.text_frame, "KEY INSIGHT", 11, TEAL, bold=True)
    tb_iv = _tb(slide, rx + 0.15, 5.42, rw - 0.25, 1.00)
    _para(tb_iv.text_frame,
          "Censoring lets us include ALL patients — not just those who experienced the event. This maximises statistical power.", 13, LIGHT)


# ── Slide 4: Kaplan-Meier Curves ──────────────────────────────────────────────

def add_km_curve_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Kaplan-Meier Curves: Visualising Survival Over Time", TEAL)

    # Chart geometry
    cl   = 1.65
    ct   = 1.25
    cw   = 9.80
    ch   = 4.70
    xmax = 60
    xs   = cw / xmax      # 0.1633 in/month
    ys   = ch             # 4.70 in per prob unit

    def t2x(t):     return cl + t * xs
    def p2y(p):     return ct + (1.0 - p) * ys

    # Chart background
    _rect(slide, cl, ct, cw, ch, CARD, "ChartBg")

    # Grid lines (horizontal at 0.25, 0.50, 0.75)
    for gp in [0.25, 0.50, 0.75]:
        gy = p2y(gp)
        _rect(slide, cl, gy, cw, 0.02, CARD_ALT, f"GridH{int(gp*100)}")

    # Y-axis label (rotated via shape.rotation)
    tb_yl = _tb(slide, 0.0, 2.80, 1.50, 1.20, "YAxisLabel")
    _para(tb_yl.text_frame, "Survival Probability", 13, MUTED, align=PP_ALIGN.CENTER)
    tb_yl.rotation = 270

    # Y-axis ticks and labels
    for gp, lbl in [(1.00, "1.00"), (0.75, "0.75"), (0.50, "0.50"), (0.25, "0.25"), (0.00, "0")]:
        gy = p2y(gp)
        _rect(slide, cl - 0.10, gy - 0.01, 0.10, 0.03, MUTED, f"YTick{int(gp*100)}")
        tb_y = _tb(slide, cl - 0.65, gy - 0.15, 0.54, 0.30)
        _para(tb_y.text_frame, lbl, 10, MUTED, align=PP_ALIGN.RIGHT)

    # X-axis line
    _rect(slide, cl, p2y(0.0), cw, 0.03, MUTED, "XAxisLine")

    # X-axis ticks and labels
    for t in [0, 12, 24, 36, 48, 60]:
        tx = t2x(t)
        _rect(slide, tx - 0.02, p2y(0.0), 0.04, 0.10, MUTED, f"XTick{t}")
        tb_t = _tb(slide, tx - 0.28, p2y(0.0) + 0.12, 0.56, 0.28)
        _para(tb_t.text_frame, str(t), 10, MUTED, align=PP_ALIGN.CENTER)

    tb_xt = _tb(slide, cl, p2y(0.0) + 0.43, cw, 0.28)
    _para(tb_xt.text_frame, "Time (months)", 13, MUTED, align=PP_ALIGN.CENTER)

    # Helper: draw a KM step function
    def _draw_km(steps, color, prefix):
        for i, (t, p) in enumerate(steps):
            t_next = steps[i + 1][0] if i + 1 < len(steps) else xmax
            p_next = steps[i + 1][1] if i + 1 < len(steps) else p
            # Horizontal segment
            x = t2x(t)
            w = t2x(t_next) - x
            y = p2y(p)
            _rect(slide, x, y, w, 0.06, color, f"{prefix}_h{i}")
            # Vertical drop
            if i + 1 < len(steps) and p_next < p:
                vy = p2y(p)
                vh = (p - p_next) * ys
                vx = t2x(t_next) - 0.03
                _rect(slide, vx, vy, 0.06, vh, color, f"{prefix}_v{i}")

    # CORAL: high expression (worse survival)
    coral_steps = [
        (0, 1.00), (6, 0.87), (12, 0.73), (18, 0.60),
        (26, 0.48), (34, 0.37), (44, 0.27), (54, 0.20), (60, 0.16),
    ]
    # TEAL: low expression (better survival)
    teal_steps = [
        (0, 1.00), (10, 0.95), (20, 0.88), (30, 0.81),
        (40, 0.75), (50, 0.69), (60, 0.64),
    ]

    _draw_km(coral_steps, CORAL, "KMH")
    _draw_km(teal_steps,  TEAL,  "KML")

    # Median survival dashed line for CORAL (p=0.50 → t≈22)
    med_y = p2y(0.50)
    _rect(slide, cl, med_y, cw * 0.37, 0.03, AMBER, "MedianH")
    _rect(slide, t2x(22) - 0.015, ct, 0.03, med_y - ct, AMBER, "MedianV")
    tb_med = _tb(slide, t2x(22) + 0.05, med_y - 0.28, 1.8, 0.28)
    _para(tb_med.text_frame, "Median ≈ 22 mo", 10, AMBER)

    # Legend (inside chart, top-right)
    leg_x = cl + cw - 2.85
    leg_y = ct + 0.20
    _rect(slide, leg_x, leg_y, 2.75, 1.00, CARD, "LegBox")
    _rect(slide, leg_x + 0.12, leg_y + 0.22, 0.42, 0.06, CORAL, "LegCoralLine")
    tb_lc = _tb(slide, leg_x + 0.60, leg_y + 0.10, 2.0, 0.30)
    _para(tb_lc.text_frame, "High expression", 11, CORAL)
    _rect(slide, leg_x + 0.12, leg_y + 0.64, 0.42, 0.06, TEAL, "LegTealLine")
    tb_lt = _tb(slide, leg_x + 0.60, leg_y + 0.52, 2.0, 0.30)
    _para(tb_lt.text_frame, "Low expression", 11, TEAL)

    # p-value badge
    _rect(slide, cl + 0.15, ct + 0.15, 2.40, 0.48, TEAL_DARK, "PvalBox")
    tb_pv = _tb(slide, cl + 0.25, ct + 0.22, 2.20, 0.35)
    _para(tb_pv.text_frame, "Log-rank  p < 0.001", 12, TEAL, bold=True)


# ── Slide 5: Hazard Ratio ──────────────────────────────────────────────────────

def add_hr_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Hazard Ratio: One Number to Rule Them All", PURPLE)

    # ── Left panel: number line ────────────────────────────────────────────────
    # Number line from HR=0.25 to HR=4.0 mapped to x=0.5 to x=6.5 (6.0" wide)
    nl_left  = 0.60
    nl_right = 6.50
    nl_w     = nl_right - nl_left
    nl_y     = 3.80
    hr_min, hr_max = 0.25, 4.0

    def hr2x(hr):
        return nl_left + (hr - hr_min) / (hr_max - hr_min) * nl_w

    # Protective zone (HR < 1)
    zone_l_w = hr2x(1.0) - nl_left
    _rect(slide, nl_left, nl_y - 0.22, zone_l_w, 0.44, TEAL_DIM, "ZoneTeal")
    # Harmful zone (HR > 1)
    zone_r_x = hr2x(1.0)
    zone_r_w = nl_right - zone_r_x
    _rect(slide, zone_r_x, nl_y - 0.22, zone_r_w, 0.44, CORAL_DIM, "ZoneCoral")
    # Main axis line
    _rect(slide, nl_left, nl_y - 0.02, nl_w, 0.04, LIGHT, "NLAxis")
    # Null vertical line at HR=1
    _rect(slide, hr2x(1.0) - 0.03, nl_y - 0.50, 0.06, 1.00, WHITE, "NullLine")

    # Tick marks and labels
    for hr, lbl in [(0.5, "0.5"), (1.0, "1.0"), (2.0, "2.0"), (4.0, "4.0")]:
        tx = hr2x(hr)
        _rect(slide, tx - 0.025, nl_y + 0.22, 0.05, 0.14, MUTED, f"NLTick{hr}")
        tb_t = _tb(slide, tx - 0.28, nl_y + 0.38, 0.56, 0.28)
        _para(tb_t.text_frame, lbl, 11, MUTED, align=PP_ALIGN.CENTER)

    tb_ax = _tb(slide, nl_left, nl_y + 0.68, nl_w, 0.28)
    _para(tb_ax.text_frame, "Hazard Ratio (HR)", 13, MUTED, align=PP_ALIGN.CENTER)

    # Zone labels
    tb_zl = _tb(slide, nl_left + 0.10, nl_y - 1.05, zone_l_w - 0.2, 0.60)
    _para(tb_zl.text_frame, "PROTECTIVE\nGene lowers risk", 13, TEAL,
          bold=True, align=PP_ALIGN.CENTER)
    tb_zr = _tb(slide, zone_r_x + 0.10, nl_y - 1.05, zone_r_w - 0.2, 0.60)
    _para(tb_zr.text_frame, "HARMFUL\nGene raises risk", 13, CORAL,
          bold=True, align=PP_ALIGN.LEFT)

    # "No effect" label at HR=1
    tb_ne = _tb(slide, hr2x(1.0) - 0.50, nl_y - 0.56, 1.00, 0.30)
    _para(tb_ne.text_frame, "HR = 1", 11, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Example annotations below the axis
    # Gene A: HR = 0.35 (protective)
    gax = hr2x(0.35)
    _rect(slide, gax - 0.03, nl_y + 0.22, 0.06, 0.20, TEAL, "GeneAMark")
    tb_ga = _tb(slide, gax - 1.0, nl_y + 1.10, 2.0, 0.55)
    _para(tb_ga.text_frame, "HR = 0.35\nGene A protects", 11, TEAL, align=PP_ALIGN.CENTER)

    # Gene B: HR = 2.8 (harmful)
    gbx = hr2x(2.8)
    _rect(slide, gbx - 0.03, nl_y + 0.22, 0.06, 0.20, CORAL, "GeneBMark")
    tb_gb = _tb(slide, gbx - 1.0, nl_y + 1.10, 2.0, 0.55)
    _para(tb_gb.text_frame, "HR = 2.8\n2.8× higher risk", 11, CORAL, align=PP_ALIGN.CENTER)

    # Interpretation cards at top-left
    _rect(slide, 0.5, 1.15, 3.0, 0.72, CARD, "InterpCard1")
    tb_i1 = _tb(slide, 0.65, 1.23, 2.70, 0.56)
    _para(tb_i1.text_frame, "HR < 1  →  gene is protective", 13, TEAL)
    _rect(slide, 0.5, 1.97, 3.0, 0.72, CARD, "InterpCard2")
    tb_i2 = _tb(slide, 0.65, 2.05, 2.70, 0.56)
    _para(tb_i2.text_frame, "HR = 1  →  no effect on survival", 13, LIGHT)
    _rect(slide, 0.5, 2.79, 3.0, 0.72, CARD, "InterpCard3")
    tb_i3 = _tb(slide, 0.65, 2.87, 2.70, 0.56)
    _para(tb_i3.text_frame, "HR > 1  →  gene increases risk", 13, CORAL)

    # ── Right panel: mini forest plot ─────────────────────────────────────────
    rx = 7.00
    chart_x = 9.30   # CI bars start here
    chart_w = 2.90   # width of CI bar area
    hr_min_f, hr_max_f = 0.5, 4.5
    f_scale = chart_w / (hr_max_f - hr_min_f)

    def fhr2x(hr):
        return chart_x + (hr - hr_min_f) * f_scale

    tb_rh = _tb(slide, rx, 1.10, 5.50, 0.30)
    _para(tb_rh.text_frame, "FOREST PLOT: EFFECT ACROSS DATASETS", 12, MUTED, bold=True)

    # Null line at HR=1.0
    null_fx = fhr2x(1.0)
    _rect(slide, null_fx - 0.02, 1.45, 0.04, 4.10, MUTED, "FNullLine")

    datasets_f = [
        ("GSE12345 (n=87)",  1.9, 1.2, 2.9),
        ("GSE23456 (n=120)", 2.2, 1.6, 3.1),
        ("GSE34567 (n=64)",  1.7, 0.9, 3.2),
        ("GSE45678 (n=210)", 2.0, 1.5, 2.7),
        ("GSE56789 (n=95)",  2.3, 1.7, 3.2),
    ]
    pooled_f = ("Pooled (n=576)", 2.05, 1.62, 2.60)

    row_ys = [1.55, 2.25, 2.95, 3.65, 4.35]
    for (lbl, hr, cil, cih), ry in zip(datasets_f, row_ys):
        tb_dl = _tb(slide, rx, ry + 0.05, chart_x - rx - 0.15, 0.40)
        _para(tb_dl.text_frame, lbl, 11, LIGHT)
        bar_x = fhr2x(cil)
        bar_w = fhr2x(cih) - bar_x
        bar_color = CORAL if cil > 1.0 else (TEAL if cih < 1.0 else MUTED)
        _rect(slide, bar_x, ry + 0.27, bar_w, 0.07, bar_color, f"FBar_{lbl[:8]}")
        dot_x = fhr2x(hr)
        _rect(slide, dot_x - 0.07, ry + 0.13, 0.14, 0.32, bar_color, f"FDot_{lbl[:8]}")

    # Pooled estimate
    py = 5.15
    pl, ph = fhr2x(pooled_f[2]), fhr2x(pooled_f[3])
    _rect(slide, rx, py + 0.05, chart_x - rx - 0.15, 0.40)
    tb_pl = _tb(slide, rx, py + 0.05, chart_x - rx - 0.15, 0.40)
    _para(tb_pl.text_frame, pooled_f[0], 11, TEAL, bold=True)
    _rect(slide, pl, py + 0.22, ph - pl, 0.30, TEAL, "FPooledBar")

    # Scale ticks
    for hr_t in [1.0, 2.0, 3.0]:
        tx_f = fhr2x(hr_t)
        _rect(slide, tx_f - 0.02, 5.60, 0.04, 0.10, MUTED, f"FTick{hr_t}")
        tb_ft = _tb(slide, tx_f - 0.20, 5.72, 0.40, 0.25)
        _para(tb_ft.text_frame, str(hr_t), 9, MUTED, align=PP_ALIGN.CENTER)


# ── Slide 6: Meta-Analysis ────────────────────────────────────────────────────

def add_meta_analysis_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Meta-Analysis: Why N Datasets > 1 Dataset", BLUE)

    # ── Left panel: ONE STUDY ─────────────────────────────────────────────────
    lw = 5.50
    _rect(slide, 0.5, 1.10, lw, 5.60, CARD, "LeftPanel")
    tb_lh = _tb(slide, 0.65, 1.22, lw - 0.20, 0.36)
    _para(tb_lh.text_frame, "ONE STUDY", 16, AMBER, bold=True)

    # n badge
    _rect(slide, 1.00, 1.75, 2.00, 0.85, CARD_ALT, "NBadge")
    tb_n = _tb(slide, 1.05, 1.80, 1.90, 0.50)
    _para(tb_n.text_frame, "n = 30", 30, AMBER, bold=True, align=PP_ALIGN.CENTER)
    tb_ns = _tb(slide, 1.05, 2.32, 1.90, 0.28)
    _para(tb_ns.text_frame, "small study", 11, MUTED, align=PP_ALIGN.CENTER)

    # Single wide CI bar (HR=2.1, CI=[0.7, 5.5])
    bar_x = 0.80
    bar_area_w = 4.90
    bar_hr_min, bar_hr_max = 0.0, 6.0
    bar_scale = bar_area_w / (bar_hr_max - bar_hr_min)

    def l2x(hr): return bar_x + hr * bar_scale

    null_lx = l2x(1.0)
    _rect(slide, null_lx - 0.02, 2.80, 0.04, 1.40, MUTED, "LNullLine")
    _rect(slide, l2x(0.7), 3.38, l2x(5.5) - l2x(0.7), 0.08, AMBER, "LCIBar")
    _rect(slide, l2x(2.1) - 0.08, 3.24, 0.16, 0.36, AMBER, "LDot")

    # CI labels
    tb_cil = _tb(slide, l2x(0.7) - 0.10, 3.50, 0.55, 0.28)
    _para(tb_cil.text_frame, "0.7", 10, MUTED)
    tb_cih = _tb(slide, l2x(5.5) - 0.10, 3.50, 0.55, 0.28)
    _para(tb_cih.text_frame, "5.5", 10, MUTED)
    tb_hr_l = _tb(slide, l2x(2.1) - 0.25, 3.60, 0.60, 0.28)
    _para(tb_hr_l.text_frame, "HR=2.1", 10, AMBER, align=PP_ALIGN.CENTER)

    # Scale labels
    for hr_v, hr_lbl in [(0, "0"), (1, "1"), (2, "2"), (4, "4"), (6, "6")]:
        _tb_v = _tb(slide, l2x(hr_v) - 0.15, 3.88, 0.30, 0.25)
        _para(_tb_v.text_frame, hr_lbl, 9, MUTED, align=PP_ALIGN.CENTER)
    _tb_ax = _tb(slide, bar_x, 4.16, bar_area_w, 0.25)
    _para(_tb_ax.text_frame, "Hazard Ratio", 10, MUTED, align=PP_ALIGN.CENTER)

    # Verdict card
    _rect(slide, 0.65, 4.50, lw - 0.30, 1.30, CARD_ALT, "VerdictL")
    tb_vl1 = _tb(slide, 0.80, 4.58, lw - 0.55, 0.36)
    _para(tb_vl1.text_frame, "p = 0.13  (not significant)", 14, AMBER, bold=True)
    tb_vl2 = _tb(slide, 0.80, 5.00, lw - 0.55, 0.72)
    _para(tb_vl2.text_frame,
          "CI is wide — crosses HR=1. Effect looks real, but statistical power is too low to be certain.", 13, LIGHT)

    # ── Right panel: 8 DATASETS ───────────────────────────────────────────────
    rp_x = 6.30
    rp_w = 12.833 - rp_x - 0.2
    _rect(slide, rp_x, 1.10, rp_w, 5.60, CARD, "RightPanel")
    tb_rh = _tb(slide, rp_x + 0.15, 1.22, rp_w - 0.25, 0.36)
    _para(tb_rh.text_frame, "8 INDEPENDENT COHORTS", 16, TEAL, bold=True)

    datasets_m = [
        ("GSE11111 (n=47)",  2.1, 1.3, 3.4),
        ("GSE22222 (n=82)",  1.9, 1.2, 3.0),
        ("GSE33333 (n=55)",  2.3, 1.5, 3.6),
        ("GSE44444 (n=130)", 2.0, 1.4, 2.9),
        ("GSE55555 (n=73)",  1.8, 1.1, 2.9),
        ("GSE66666 (n=96)",  2.2, 1.5, 3.2),
        ("GSE77777 (n=61)",  2.4, 1.6, 3.7),
        ("GSE88888 (n=108)", 1.9, 1.3, 2.8),
    ]
    pooled_m = ("Pooled (n=652)", 2.04, 1.68, 2.48)

    mc_x   = rp_x + 2.30   # CI bars start
    mc_w   = rp_x + rp_w - mc_x - 0.15
    mc_min, mc_max = 0.5, 4.5
    mc_scale = mc_w / (mc_max - mc_min)

    def m2x(hr): return mc_x + (hr - mc_min) * mc_scale

    null_mx = m2x(1.0)
    _rect(slide, null_mx - 0.02, 1.65, 0.04, 3.50, MUTED, "MNullLine")

    row_ys_m = [1.68, 2.08, 2.48, 2.88, 3.28, 3.68, 4.08, 4.48]
    for (lbl, hr, cil, cih), ry in zip(datasets_m, row_ys_m):
        tb_ml = _tb(slide, rp_x + 0.10, ry + 0.04, 2.10, 0.30)
        _para(tb_ml.text_frame, lbl, 9, MUTED)
        bx  = m2x(cil)
        bw  = m2x(cih) - bx
        _rect(slide, bx, ry + 0.13, bw, 0.05, TEAL, f"MBar{lbl[:5]}")
        _rect(slide, m2x(hr) - 0.05, ry + 0.06, 0.10, 0.22, TEAL, f"MDot{lbl[:5]}")

    # Pooled bar
    py_m = 5.00
    _rect(slide, rp_x + 0.10, py_m + 0.04, 2.10, 0.35)
    tb_pm = _tb(slide, rp_x + 0.10, py_m + 0.04, 2.10, 0.35)
    _para(tb_pm.text_frame, pooled_m[0], 10, TEAL, bold=True)
    px1, px2 = m2x(pooled_m[2]), m2x(pooled_m[3])
    _rect(slide, px1, py_m + 0.08, px2 - px1, 0.22, TEAL, "MPooled")

    # Stats
    tb_stats = _tb(slide, rp_x + 0.10, 5.45, rp_w - 0.20, 0.68)
    _para(tb_stats.text_frame,
          "p_pooled < 0.001  ·  I² = 23% (low heterogeneity)  ·  Consistent across all 8 cohorts", 12, TEAL, bold=True)

    # Bottom analogy badge
    _rect(slide, 0.50, 6.82, 12.333, 0.52, TEAL_DARK, "AnalogyBadge")
    tb_an = _tb(slide, 0.65, 6.88, 12.0, 0.38)
    _para(tb_an.text_frame,
          "For software engineers: meta-analysis is cross-validation across independent labs — reproducibility over sample size",
          15, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── Slide 7: What We Built ─────────────────────────────────────────────────────

def add_comparison_table_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "What We Built: NL Query → Results in 5 Minutes", TEAL)

    tb_tag = _tb(slide, 0.5, 1.02, 12.333, 0.32)
    _para(tb_tag.text_frame,
          "Type a question in plain English → cross-cohort survival analysis → publication-ready export in under 5 minutes",
          14, MUTED)

    col_headers = ["Feature", "KMplot", "GEPIA2", "OncoLnc", "GEO Survival Analysis  ✓"]
    col_widths  = [3.10, 1.80, 1.80, 1.80, 3.853]
    rows = [
        ["Data source",        "TCGA + curated",  "TCGA + GTEx",   "TCGA only",    "All of GEO  (200K+ datasets)"],
        ["Custom datasets",    "✗",               "✗",             "✗",            "✓  Any GSE study"],
        ["Cross-cohort meta",  "✗",               "✗",             "✗",            "✓  N-dataset pooling"],
        ["Natural language",   "✗",               "✗",             "✗",            "✓  Plain English query"],
        ["Forest plots + I²",  "✗",               "Partial",       "✗",            "✓  + heterogeneity stats"],
        ["Open source / free", "✗",               "✗",             "✗",            "✓  MIT license"],
    ]

    n_cols = len(col_headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.38),
        Inches(12.333), Inches(5.80),
    ).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ci, hdr in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        tf = cell.text_frame
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        cell.fill.solid()
        if ci == 4:
            cell.fill.fore_color.rgb = TEAL
            tf.paragraphs[0].font.color.rgb = BG
        else:
            cell.fill.fore_color.rgb = CARD
            tf.paragraphs[0].font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        fill = CARD if ri % 2 == 0 else CARD_ALT
        for ci, text in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = text
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(13)
            cell.fill.solid()
            if ci == 4:
                cell.fill.fore_color.rgb = TEAL_DARK
                tf.paragraphs[0].font.color.rgb = TEAL
                tf.paragraphs[0].font.bold = True
            elif ci == 0:
                cell.fill.fore_color.rgb = fill
                tf.paragraphs[0].font.color.rgb = WHITE
                tf.paragraphs[0].font.bold = True
            else:
                cell.fill.fore_color.rgb = fill
                if text == "✗":
                    tf.paragraphs[0].font.color.rgb = CORAL
                elif text == "✓":
                    tf.paragraphs[0].font.color.rgb = TEAL
                else:
                    tf.paragraphs[0].font.color.rgb = LIGHT


# ── Slide 8 & 9: Architecture slides ──────────────────────────────────────────

def add_architecture_slide(
    prs: Presentation,
    title: str,
    layers: list[tuple[str, list[str], RGBColor]],
    accent: RGBColor = TEAL,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, title, accent)

    n_layers  = len(layers)
    row_h     = 0.88
    arrow_h   = 0.14
    total     = n_layers * row_h + (n_layers - 1) * arrow_h
    start_y   = 1.05 + (6.30 - total) / 2

    label_x    = 0.30
    label_w    = 1.50
    box_area_x = 2.00
    box_area_w = 11.133
    gap        = 0.10

    for i, (label, boxes, color) in enumerate(layers):
        y = start_y + i * (row_h + arrow_h)
        tb_lbl = _tb(slide, label_x, y + 0.25, label_w, 0.40)
        _para(tb_lbl.text_frame, label, 11, MUTED, align=PP_ALIGN.RIGHT)
        n  = len(boxes)
        bw = (box_area_w - (n - 1) * gap) / n
        for j, text in enumerate(boxes):
            bx = box_area_x + j * (bw + gap)
            _rect(slide, bx, y, bw, row_h, CARD, f"ArchBox{i}_{j}")
            _rect(slide, bx, y, bw, 0.07, color, f"ArchStripe{i}_{j}")
            tb_b = _tb(slide, bx, y + 0.18, bw, row_h - 0.25)
            _para(tb_b.text_frame, text, 13, LIGHT, align=PP_ALIGN.CENTER)
        if i < n_layers - 1:
            ax = box_area_x + box_area_w / 2 - 0.06
            _rect(slide, ax, y + row_h, 0.12, arrow_h, accent, f"ArchArrow{i}")


# ── Backend / Frontend visual architecture slides ──────────────────────────────

def add_backend_arch_visual_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Backend Architecture", TEAL)

    route_x, route_w = 0.35, 3.25
    svc_x,   svc_w   = 4.00, 9.00
    row_h, gap = 0.76, 0.10
    start_y = 1.20

    routes = [
        ("GET  /api/search/stream",  "SSE analysis stream",    TEAL),
        ("POST /api/chat/…",         "Streaming AI chat",      PURPLE),
        ("POST /api/auth/…",         "JWT auth endpoints",     AMBER),
        ("POST /api/compare/…",      "Side-by-side analysis",  BLUE),
        ("POST /api/enrichment/…",   "GO/KEGG enrichment",     CORAL),
    ]
    services = [
        ("GEOSurvivalWorkflowOrchestrator",
         "chains GEOClient → GeneMapping → SurvivalAnalysis  via asyncio.gather(return_exceptions=True)",
         TEAL),
        ("PydanticAIService  ← pydantic-ai Agent",
         "5 registered tools · numpy cosine RAG · domain score badge per response",
         PURPLE),
        ("AuthService  → JWT + bcrypt",
         "register · login · /me · Depends(get_current_user) per protected route",
         AMBER),
        ("CompareService",
         "side-by-side result diff · shared result permalinks  (/results/:id  public)",
         BLUE),
        ("PathwayEnrichmentService",
         "GO term enrichment · KEGG pathways  (stub, roadmap F11)",
         CORAL),
    ]

    for i, ((route_label, route_sub, color), (svc_label, svc_desc, _)) in enumerate(
        zip(routes, services)
    ):
        ry = start_y + i * (row_h + gap)
        # Route box
        _rect(slide, route_x, ry, route_w, row_h, CARD, f"beR{i}")
        _rect(slide, route_x, ry, 0.06, row_h, color, f"beRA{i}")
        tb_r = _tb(slide, route_x + 0.14, ry + 0.08, route_w - 0.18, 0.32)
        _para(tb_r.text_frame, route_label, 10, color, bold=True, mono=True)
        tb_rs = _tb(slide, route_x + 0.14, ry + 0.42, route_w - 0.18, 0.26)
        _para(tb_rs.text_frame, route_sub, 9, MUTED)
        # Arrow
        arrow_cy = ry + row_h / 2 - 0.05
        _rect(slide, route_x + route_w + 0.06, arrow_cy, 0.42, 0.10, color, f"beAr{i}")
        _rect(slide, route_x + route_w + 0.38, arrow_cy - 0.05, 0.12, 0.20, color, f"beArH{i}")
        # Service box
        _rect(slide, svc_x, ry, svc_w, row_h, CARD, f"beS{i}")
        _rect(slide, svc_x, ry, 0.06, row_h, color, f"beSA{i}")
        tb_sh = _tb(slide, svc_x + 0.14, ry + 0.06, svc_w - 0.20, 0.32)
        _para(tb_sh.text_frame, svc_label, 12, color, bold=True)
        tb_sd = _tb(slide, svc_x + 0.14, ry + 0.40, svc_w - 0.20, 0.30)
        _para(tb_sd.text_frame, svc_desc, 9, MUTED)

    # Storage strip
    total_content_w = svc_x + svc_w - route_x
    storage_y = start_y + 5 * (row_h + gap) - gap + 0.14
    _rect(slide, route_x, storage_y, total_content_w, 0.24, TEAL_DARK, "beStorHdr")
    tb_sh = _tb(slide, route_x + 0.14, storage_y + 0.04, 2.0, 0.16)
    _para(tb_sh.text_frame, "PERSISTENT STORAGE", 10, TEAL, bold=True)

    sb_y = storage_y + 0.28
    sb_h = 7.20 - sb_y
    sb_w = (total_content_w - 2 * 0.10) / 3
    storage_items = [
        ("SQLite / PostgreSQL", "auth · chat history · analysis results\nenv-var switchable", BLUE),
        ("Parquet cache", "expression matrices · gene mappings\nprobe → gene symbol  (per-GPL)", TEAL),
        ("numpy RAG index", "Mistral 1024-dim embeddings\n500 doc chunks · ~1ms lookup", PURPLE),
    ]
    for j, (label, desc, color) in enumerate(storage_items):
        sx = route_x + j * (sb_w + 0.10)
        _rect(slide, sx, sb_y, sb_w, sb_h, CARD_ALT, f"beSto{j}")
        _rect(slide, sx, sb_y, sb_w, 0.06, color, f"beStorS{j}")
        tb_sl = _tb(slide, sx + 0.10, sb_y + 0.10, sb_w - 0.16, 0.26)
        _para(tb_sl.text_frame, label, 11, color, bold=True)
        tb_sd = _tb(slide, sx + 0.10, sb_y + 0.38, sb_w - 0.16, sb_h - 0.44)
        _para(tb_sd.text_frame, desc, 9, MUTED)


def add_frontend_arch_visual_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Frontend Architecture", BLUE)

    divx = 7.10
    lx, lw = 0.35, divx - 0.35 - 0.20
    rx, rw = divx + 0.20, 13.133 - divx - 0.20

    _rect(slide, divx, 1.10, 0.04, 6.10, BLUE, "feDiv")

    # Left: component tree as monospace text
    tb_lh = _tb(slide, lx, 1.12, lw, 0.28)
    _para(tb_lh.text_frame, "COMPONENT TREE", 11, BLUE, bold=True)

    tree_lines = [
        ("React 18 + TypeScript  (Vite + Tailwind)", LIGHT, False),
        ("│", MUTED, True),
        ("├─ <App>  (React Router v6)", LIGHT, True),
        ("│   ├─ <ChatContainer>", TEAL, True),
        ("│   │   ├─ <MessageList>  +  <ChatInput>", MUTED, True),
        ("│   │   └─ <AnalysisResults>", MUTED, True),
        ("│   │       tabs: volcano / KM / forest plot", MUTED, True),
        ("│   ├─ <AnalysisHistoryPage>", BLUE, True),
        ("│   ├─ <ComparisonPage>", BLUE, True),
        ("│   └─ <SharedResultPage>  (public, no login)", MUTED, True),
        ("│", MUTED, True),
        ("├─ Redux Toolkit Store", PURPLE, True),
        ("│   ├─ authSlice   user · token · isAuthenticated", MUTED, True),
        ("│   └─ chatSlice   messages · results · streaming", MUTED, True),
        ("│", MUTED, True),
        ("└─ API Services", AMBER, True),
        ("    ├─ api.ts       EventSource SSE  →  /search/stream", MUTED, True),
        ("    ├─ chatApi.ts   fetch stream     →  /chat/…", MUTED, True),
        ("    └─ authApi.ts   axios            →  /auth/token", MUTED, True),
    ]

    tb_tree = _tb(slide, lx, 1.48, lw, 5.60)
    tf = tb_tree.text_frame
    tf.word_wrap = False
    for k, (line, color, is_new) in enumerate(tree_lines):
        if k == 0:
            _para(tf, line, 10, color, mono=True)
        else:
            _para(tf, line, 10, color, new=True, mono=True)

    # Right: 3 info cards
    tb_rh = _tb(slide, rx, 1.12, rw, 0.28)
    _para(tb_rh.text_frame, "STATE + DATA FLOW", 11, BLUE, bold=True)

    _ic(slide, rx, 1.48, rw, 1.80, "REDUX STATE FLOW", PURPLE, [
        "User submits query  →  dispatch(sendMessage())",
        "Thunk calls EventSource /api/search/stream",
        "onmessage  →  dispatch(setAnalysisProgress(event))",
        "<AnalysisResults> re-renders on each SSE event",
        "On complete  →  dispatch(setResults(data))  →  tabs render",
    ], "feRx")

    _ic(slide, rx, 3.42, rw, 1.55, "TWO STREAMING PROTOCOLS", TEAL, [
        "Analysis: EventSource SSE — auto-reconnect, browser-native",
        "Chat: fetch + ReadableStream — manual token-by-token append",
        "Both dispatch to chatSlice  →  one reactive UI pattern",
        "SSE keeps connection open; chat opens per-message",
    ], "feSSE")

    _ic(slide, rx, 5.10, rw, 2.10, "WHY REDUX TOOLKIT?", BLUE, [
        "authSlice: single source of truth for auth state across all routes",
        "chatSlice: streaming flag prevents double-send on re-render",
        "time-travel DevTools: debug state transitions step-by-step",
        "condition callbacks: fetchConversations skips if data <30s old",
        "No prop-drilling: deep components read store directly via useSelector",
    ], "feRdx")


# ── Slide 10: Workflow ─────────────────────────────────────────────────────────

def add_workflow_slide(
    prs: Presentation,
    title: str,
    steps: list[tuple[str, str, str]],
    bottom_text: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, title, CORAL)

    row_defs = [(0, 3), (3, 6)]
    row_tops = [1.10, 3.85]
    col_xs   = [0.5, 4.62, 8.73]
    sw, sh   = 3.90, 2.60

    for row_i, (start, end) in enumerate(row_defs):
        row_steps = steps[start:end]
        yt = row_tops[row_i]
        for col_i, (num, label, desc) in enumerate(row_steps):
            x = col_xs[col_i]
            _rect(slide, x, yt, sw, sh, CARD, f"Step{start + col_i}")
            tb_n = _tb(slide, x + 0.15, yt + 0.10, sw - 0.30, 0.40)
            _para(tb_n.text_frame, num, 16, CORAL, bold=True, align=PP_ALIGN.CENTER)
            tb_l = _tb(slide, x + 0.15, yt + 0.60, sw - 0.30, 0.45)
            _para(tb_l.text_frame, label, 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
            tb_d = _tb(slide, x + 0.15, yt + 1.15, sw - 0.30, 1.30)
            _para(tb_d.text_frame, desc, 13, LIGHT, align=PP_ALIGN.CENTER)
            if col_i < len(row_steps) - 1:
                ax = x + sw + 0.02
                ay = yt + sh / 2 - 0.04
                _rect(slide, ax, ay, 0.20, 0.08, CORAL, f"Arrow{start+col_i}")
        if row_i == 0:
            _rect(slide, col_xs[0] + sw / 2 - 0.04, 3.65, 0.08, 0.20, CORAL, "DownArrow")

    if bottom_text:
        _rect(slide, 0.5, 6.60, 12.333, 0.65, TEAL, "BottomHighlight")
        tb_bt = _tb(slide, 0.5, 6.63, 12.333, 0.55)
        _para(tb_bt.text_frame, bottom_text, 16, BG, bold=True, align=PP_ALIGN.CENTER)


# ── Slide 11: Gene Mapping Babel ──────────────────────────────────────────────

def add_gene_babel_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Gene Mapping: The Platform Babel Problem", CORAL)

    # ── Left: problem ─────────────────────────────────────────────────────────
    tb_lh = _tb(slide, 0.5, 1.10, 5.80, 0.30)
    _para(tb_lh.text_frame, "RAW PROBE IDs — NOT GENE NAMES", 12, CORAL, bold=True)

    probe_data = [
        ("Affymetrix HG-U133+",  '"1007_s_at"',    '"1053_at"',    '"117_at"'),
        ("Illumina HumanHT-12",  '"ILMN_1343291"', '"ILMN_1343295"', '"ILMN_1651209"'),
        ("Agilent SurePrint",    '"A_23_P100001"', '"A_23_P100011"', '"A_23_P100022"'),
    ]
    card_top = 1.50
    for idx, (platform, p1, p2, p3) in enumerate(probe_data):
        cy = card_top + idx * 1.52
        _rect(slide, 0.5, cy, 5.60, 1.38, CARD, f"ProbeCard{idx}")
        _rect(slide, 0.5, cy, 0.08, 1.38, CORAL, f"ProbeAccent{idx}")
        tb_plat = _tb(slide, 0.68, cy + 0.08, 5.30, 0.30)
        _para(tb_plat.text_frame, platform, 12, CORAL, bold=True)
        tb_ids = _tb(slide, 0.68, cy + 0.44, 5.30, 0.76)
        _para(tb_ids.text_frame, f"{p1}   {p2}   {p3}", 12, AMBER)
        # Arrow pointing to ???
        _rect(slide, 5.90, cy + 0.55, 0.40, 0.06, CORAL, f"ProbeArrow{idx}")
        tb_q = _tb(slide, 6.32, cy + 0.40, 0.60, 0.35)
        _para(tb_q.text_frame, "???", 18, CORAL, bold=True)

    tb_stat = _tb(slide, 0.5, 6.05, 5.80, 0.50)
    _para(tb_stat.text_frame, "100+ platforms · Each with proprietary IDs · No universal standard",
          13, MUTED, align=PP_ALIGN.CENTER)

    # ── Divider ────────────────────────────────────────────────────────────────
    _rect(slide, 7.0, 1.05, 0.04, 5.80, TEAL, "Divider")

    # ── Right: solution ───────────────────────────────────────────────────────
    rx = 7.20
    rw = 12.633 - rx

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "OUR AUTOMATED SOLUTION", 12, TEAL, bold=True)

    steps_r = [
        (BLUE,   "GSE arrives with expression matrix"),
        (TEAL,   "GPL annotation file auto-downloaded"),
        (TEAL,   "probe_id → gene_symbol column detected"),
        (PURPLE, "Multi-probe genes aggregated  (mean)"),
        (AMBER,  "Non-expression platforms filtered out"),
        (TEAL,   "Cached to disk  (platform_mappings/*.parquet)"),
    ]
    step_top = 1.50
    step_h   = 0.72
    arrow_h  = 0.12
    for i, (color, text) in enumerate(steps_r):
        sy = step_top + i * (step_h + arrow_h)
        _rect(slide, rx, sy, rw, step_h, CARD, f"SolStep{i}")
        _rect(slide, rx, sy, 0.08, step_h, color, f"SolAccent{i}")
        tb_s = _tb(slide, rx + 0.18, sy + 0.18, rw - 0.28, step_h - 0.22)
        _para(tb_s.text_frame, text, 13, LIGHT)
        if i < len(steps_r) - 1:
            _rect(slide, rx + rw / 2 - 0.04, sy + step_h, 0.08, arrow_h, TEAL, f"SolArrow{i}")

    _rect(slide, rx, 6.05, rw, 0.50, TEAL_DARK, "SolBadge")
    tb_sb = _tb(slide, rx + 0.10, 6.13, rw - 0.15, 0.34)
    _para(tb_sb.text_frame, "Zero manual curation  ·  Any platform  ·  Any study",
          13, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── Slide 12: Gene Mapping Cache ──────────────────────────────────────────────

def add_gene_cache_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Gene Mapping: Caching Makes It Possible", TEAL)

    # ── Left: 3-tier cache funnel ─────────────────────────────────────────────
    tb_lh = _tb(slide, 0.5, 1.10, 5.30, 0.30)
    _para(tb_lh.text_frame, "3-TIER CACHE STRATEGY", 12, TEAL, bold=True)

    tiers = [
        # (left_offset, width, color, title, lines, hit_color)
        (0.80, 4.30, TEAL,   "In-Memory LRU",     ["Max 3 platforms", "Hit time: 0 ms", "Cleared on restart"], TEAL),
        (0.50, 4.90, BLUE,   "Disk Parquet Cache", ["101 platforms cached", "Hit time: ~50 ms", "Persists across restarts"], BLUE),
        (0.20, 5.50, AMBER,  "NCBI FTP Download",  ["~1 GB avg per platform", "Download: 2–30 min", "Triggered only on cache miss"], AMBER),
    ]
    tier_top = 1.50
    tier_h   = 1.30
    tier_gap = 0.18
    for idx, (lo, tw, color, title, lines, hc) in enumerate(tiers):
        ty = tier_top + idx * (tier_h + tier_gap)
        _rect(slide, lo, ty, tw, tier_h, CARD, f"Tier{idx}")
        _rect(slide, lo, ty, tw, 0.07, color, f"TierTop{idx}")
        tb_tt = _tb(slide, lo + 0.15, ty + 0.12, tw - 0.25, 0.36)
        _para(tb_tt.text_frame, title, 15, color, bold=True)
        for li, line in enumerate(lines):
            tb_l = _tb(slide, lo + 0.15, ty + 0.52 + li * 0.26, tw - 0.25, 0.26)
            _para(tb_l.text_frame, f"• {line}", 12, MUTED)
        if idx < len(tiers) - 1:
            miss_y = ty + tier_h + 0.02
            _rect(slide, 2.30, miss_y, 1.50, tier_gap - 0.04, CARD_ALT, f"MissBox{idx}")
            tb_m = _tb(slide, 2.35, miss_y + 0.01, 1.40, tier_gap - 0.06)
            _para(tb_m.text_frame, "MISS →", 10, CORAL, bold=True, align=PP_ALIGN.CENTER)

    # ── Right: power-law bar chart ────────────────────────────────────────────
    rx = 6.30
    rw = 12.633 - rx
    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "PLATFORM USAGE IN GEO (POWER LAW)", 12, TEAL, bold=True)

    # Approximate % of human cancer expression datasets per platform
    platforms_pct = [
        ("GPL570",   33),
        ("GPL96",    13),
        ("GPL10558",  9),
        ("GPL6947",   8),
        ("GPL6244",   6),
        ("GPL13667",  5),
        ("GPL6480",   4),
        ("GPL17077",  3),
        ("Others",   19),
    ]
    n_bars   = len(platforms_pct)
    bar_area_left  = rx + 0.10
    bar_area_top   = 1.55
    bar_area_h     = 4.00
    bar_area_w     = rw - 0.25
    slot_w         = bar_area_w / n_bars
    bar_w          = slot_w * 0.68
    max_pct        = 33

    bar_bottom = bar_area_top + bar_area_h
    # Axis
    _rect(slide, bar_area_left, bar_bottom, bar_area_w, 0.03, MUTED, "BarAxis")

    for i, (platform, pct) in enumerate(platforms_pct):
        bh     = (pct / max_pct) * bar_area_h
        bx     = bar_area_left + i * slot_w + (slot_w - bar_w) / 2
        by     = bar_bottom - bh
        color  = TEAL if platform != "Others" else MUTED
        _rect(slide, bx, by, bar_w, bh, color, f"BarP{i}")
        # Percentage label above bar
        tb_pv = _tb(slide, bx - 0.05, by - 0.28, bar_w + 0.10, 0.26)
        _para(tb_pv.text_frame, f"{pct}%", 10, color, bold=True, align=PP_ALIGN.CENTER)
        # Platform label below axis (rotated look via small text)
        tb_pl = _tb(slide, bx - 0.05, bar_bottom + 0.05, bar_w + 0.10, 0.35)
        _para(tb_pl.text_frame, platform, 8, MUTED, align=PP_ALIGN.CENTER)

    # Cumulative annotation
    _rect(slide, bar_area_left, 6.07, bar_area_w, 0.52, TEAL_DARK, "CumBadge")
    tb_cum = _tb(slide, bar_area_left + 0.10, 6.13, bar_area_w - 0.15, 0.38)
    _para(tb_cum.text_frame,
          "Top 8 platforms ≈ 81% of datasets  →  cache them first, prevent 80% of 2–30 min downloads",
          12, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── Slide 13: Statistical Engine ──────────────────────────────────────────────

def add_two_col_slide(
    prs: Presentation,
    title: str,
    left_header: str,
    left_header_color: RGBColor,
    left_bullets: list[str],
    right_header: str,
    right_bullets: list[str],
    accent: RGBColor = TEAL,
    bottom_badge: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, title, accent)

    tb_lh = _tb(slide, 0.5, 1.20, 6.0, 0.50, "LeftHeader")
    _para(tb_lh.text_frame, left_header, 18, left_header_color, bold=True)
    tb_lb = _tb(slide, 0.5, 1.80, 6.0, 4.0, "LeftBullets")
    tf = tb_lb.text_frame
    for i, bullet in enumerate(left_bullets):
        _para(tf, f"• {bullet}", 16, LIGHT, new=(i > 0), space_after=8)

    _rect(slide, 7.0, 1.20, 6.0, 4.60, CARD, "RightCard")
    tb_rh = _tb(slide, 7.15, 1.30, 5.70, 0.50, "RightHeader")
    _para(tb_rh.text_frame, right_header, 16, accent, bold=True, align=PP_ALIGN.CENTER)
    tb_rb = _tb(slide, 7.15, 1.90, 5.70, 3.80, "RightBullets")
    tf2 = tb_rb.text_frame
    for i, bullet in enumerate(right_bullets):
        _para(tf2, f"• {bullet}", 15, LIGHT, new=(i > 0), space_after=8)

    if bottom_badge:
        tb_b = _tb(slide, 0.5, 6.05, 12.333, 0.60, "BottomBadge")
        _para(tb_b.text_frame, bottom_badge, 16, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── Slide 15: SQLite vs PostgreSQL ────────────────────────────────────────────

def add_db_comparison_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "SQLite vs PostgreSQL: Cost vs Capability", BLUE)

    tb_tag = _tb(slide, 0.5, 1.02, 12.333, 0.32)
    _para(tb_tag.text_frame,
          "The app supports both via SQLAlchemy async — switch with a single environment variable",
          14, MUTED)

    col_headers = ["Feature", "SQLite  (default)", "PostgreSQL  (optional)"]
    col_widths  = [3.50, 4.20, 4.633]
    rows = [
        ["Cost",                 "$0  (file-based, zero infra)",       "$20–50/mo managed  OR  self-hosted ops"],
        ["Setup",                "Zero config · runs in 1 line",       "Docker / RDS · connection pooling required"],
        ["Concurrent writes",    "Single writer (WAL mode helps)",     "Full MVCC · unlimited concurrent writes"],
        ["Read performance",     "Excellent for < 100 concurrent",     "Scales to thousands of concurrent readers"],
        ["Migrations",           "Alembic works perfectly",            "Alembic works perfectly"],
        ["Our current use",      "Auth + results history  (low vol.)", "N/A — not needed yet"],
        ["Switch trigger",       "set  DATABASE_URL=postgresql+asyncpg://…", "Already wired · no code change needed"],
    ]

    n_cols = len(col_headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.38),
        Inches(12.333), Inches(5.90),
    ).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ci, hdr in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        tf = cell.text_frame
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        cell.fill.solid()
        header_colors = [CARD, TEAL, BLUE]
        text_colors   = [WHITE, BG, BG]
        cell.fill.fore_color.rgb = header_colors[ci]
        tf.paragraphs[0].font.color.rgb = text_colors[ci]

    for ri, row in enumerate(rows):
        fill = CARD if ri % 2 == 0 else CARD_ALT
        for ci, text in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = text
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            if ci == 0:
                tf.paragraphs[0].font.color.rgb = WHITE
                tf.paragraphs[0].font.bold = True
            elif ci == 1:
                tf.paragraphs[0].font.color.rgb = TEAL
            else:
                tf.paragraphs[0].font.color.rgb = BLUE


# ── Deployment / gap-cards / image slides ─────────────────────────────────────

def add_gap_cards_slide(
    prs: Presentation,
    title: str,
    cards: list[tuple[str, str, RGBColor]],
    accent: RGBColor = CORAL,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, title, accent)

    xs   = [0.5, 4.70, 8.90]
    ys   = [1.10, 3.00, 4.90]
    cw, ch = 3.90, 1.70

    for idx, (heading, desc, color) in enumerate(cards[:6]):
        col = idx % 3
        row = idx // 3
        x, y = xs[col], ys[row]
        _rect(slide, x, y, cw, ch, CARD, f"GapCard{idx}")
        tb_h = _tb(slide, x + 0.15, y + 0.10, cw - 0.30, 0.50)
        _para(tb_h.text_frame, heading, 18, color, bold=True)
        tb_d = _tb(slide, x + 0.15, y + 0.65, cw - 0.30, 0.90)
        _para(tb_d.text_frame, desc, 14, MUTED)


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tb = _tb(slide, 0.5, 0.30, 12.333, 0.75, "Title")
    _para(tb.text_frame, title, 24, WHITE, bold=True)
    _rect(slide, 0.5, 1.05, 12.333, 0.04, TEAL, "TitleLine")
    slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.20), width=Inches(11.333))
    tb_cap = _tb(slide, 0.5, 6.50, 12.333, 0.75, "Caption")
    _para(tb_cap.text_frame, caption, 14, TEAL, align=PP_ALIGN.CENTER)


def _find_screenshot(screenshots_dir: Path, *names: str) -> Path:
    subdirs = sorted(
        [d for d in screenshots_dir.iterdir() if d.is_dir()],
        reverse=True,
    )
    for name in names:
        for d in subdirs + [screenshots_dir]:
            p = d / name
            if p.exists():
                return p
    raise FileNotFoundError(f"Screenshot not found: {names} in {screenshots_dir}")


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers for deep-dive slides
# ═══════════════════════════════════════════════════════════════════════════════

def _ic(slide, x, y, w, h, header, hcol, bullets, pfx=""):
    """Info card: left accent stripe, bold header, bullet list."""
    _rect(slide, x, y, w, h, CARD, f"{pfx}IC")
    _rect(slide, x, y, 0.07, h, hcol, f"{pfx}ICA")
    tb_h = _tb(slide, x + 0.16, y + 0.08, w - 0.24, 0.30)
    _para(tb_h.text_frame, header, 11, hcol, bold=True)
    tb_b = _tb(slide, x + 0.16, y + 0.42, w - 0.24, h - 0.50)
    tf = tb_b.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        _para(tf, f"• {b}", 10, LIGHT, new=(i > 0))


def _fc(slide, items, x, y, bw, bh, gap=0.16, pfx=""):
    """Vertical flow: boxes with down-arrows between them.  items=(title, sub, color)"""
    for i, (title, sub, col) in enumerate(items):
        by = y + i * (bh + gap)
        _rect(slide, x, by, bw, bh, CARD, f"{pfx}FB{i}")
        _rect(slide, x, by, 0.07, bh, col, f"{pfx}FBA{i}")
        tb_t = _tb(slide, x + 0.14, by + 0.08, bw - 0.22, 0.32)
        _para(tb_t.text_frame, title, 12, col, bold=True)
        if sub:
            tb_s = _tb(slide, x + 0.14, by + 0.44, bw - 0.22, bh - 0.52)
            _para(tb_s.text_frame, sub, 10, MUTED)
        if i < len(items) - 1:
            ay = by + bh + 0.01
            _rect(slide, x + bw / 2 - 0.04, ay, 0.08, gap - 0.02, TEAL, f"{pfx}FA{i}")


def _fh(slide, items, x, y, bw, bh, gap=0.20, pfx="", acol=None):
    """Horizontal flow: boxes with right-arrows.  items=(title, sub, color)"""
    if acol is None:
        acol = TEAL
    for i, (title, sub, col) in enumerate(items):
        bx = x + i * (bw + gap)
        _rect(slide, bx, y, bw, bh, CARD, f"{pfx}HB{i}")
        _rect(slide, bx, y, bw, 0.06, col, f"{pfx}HBT{i}")
        tb_t = _tb(slide, bx + 0.08, y + 0.12, bw - 0.16, 0.32)
        _para(tb_t.text_frame, title, 11, col, bold=True, align=PP_ALIGN.CENTER)
        if sub:
            tb_s = _tb(slide, bx + 0.08, y + 0.48, bw - 0.16, bh - 0.56)
            _para(tb_s.text_frame, sub, 9, MUTED, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            ax = bx + bw + 0.02
            _rect(slide, ax, y + bh / 2 - 0.03, gap - 0.04, 0.06, acol, f"{pfx}HA{i}")


def _divline(slide, x=6.35, color=None):
    if color is None:
        color = TEAL
    _rect(slide, x, 1.05, 0.04, 5.80, color, "Divider")


# ── N1: React 18 ──────────────────────────────────────────────────────────────

def add_react_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "React 18: Component-Based UI Framework", BLUE)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS REACT?", BLUE, [
        "UI library — every piece of UI is a TypeScript function returning JSX",
        "Virtual DOM: computes minimal real DOM changes for performance",
        "Hooks: useState / useEffect / useCallback — stateful logic in functions",
        "Unidirectional data flow: parent passes props down to children",
    ], "ra1")
    _ic(slide, 0.5, 2.72, 5.80, 1.50, "WHY REACT?", BLUE, [
        "TypeScript-native: all props and API shapes are compile-time typed",
        "Redux Toolkit hooks (useSelector / useDispatch) integrate natively",
        "Component reuse: AnalysisResults + KMChart rendered in 3+ pages",
        "Vite + SWC: hot module replacement — sub-second feedback in dev",
    ], "ra2")
    _ic(slide, 0.5, 4.35, 5.80, 1.50, "KEY PATTERNS IN OUR CODEBASE", BLUE, [
        "Controlled inputs: ChatInput.value ↔ Redux chatSlice state",
        "Memoization: useMemo for KM data transforms to avoid re-renders",
        "Conditional rendering: spinner ↔ result grid based on isStreaming",
        "useCallback: SSE event handler memoised to prevent re-subscriptions",
    ], "ra3")

    _divline(slide, color=BLUE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "COMPONENT TREE", 12, BLUE, bold=True)

    _rect(slide, rx + 1.0, 1.48, 4.0, 0.48, CARD, "raApp")
    _rect(slide, rx + 1.0, 1.48, 4.0, 0.06, BLUE, "raAppT")
    tb_app = _tb(slide, rx + 1.0, 1.60, 4.0, 0.28)
    _para(tb_app.text_frame, "App.tsx  (Router + AuthProvider)", 12, BLUE, bold=True, align=PP_ALIGN.CENTER)
    _rect(slide, rx + 3.0, 1.96, 0.06, 0.20, BLUE, "raA0")

    pages = [("ChatContainer", TEAL), ("HistoryPage", PURPLE), ("ComparisonPage", AMBER), ("SharedResult", CORAL)]
    pw = (rw - 0.45) / 4
    for i, (name, col) in enumerate(pages):
        px = rx + i * (pw + 0.15)
        _rect(slide, px, 2.16, pw, 0.48, CARD, f"raP{i}")
        _rect(slide, px, 2.16, pw, 0.06, col, f"raPT{i}")
        tb_p = _tb(slide, px + 0.04, 2.28, pw - 0.08, 0.28)
        _para(tb_p.text_frame, name, 9, col, align=PP_ALIGN.CENTER)

    _rect(slide, rx + pw / 2 + 0.1, 2.64, 0.06, 0.18, TEAL, "raA1")

    kids = [("ChatInput", "Message box\n+ settings"), ("MessageList", "Chat history\n+ streaming"), ("AnalysisResults", "KM curves\n+ gene table")]
    kw = (rw - 0.30) / 3
    for i, (name, desc) in enumerate(kids):
        kx = rx + i * (kw + 0.15)
        _rect(slide, kx, 2.82, kw, 0.80, CARD, f"raK{i}")
        _rect(slide, kx, 2.82, kw, 0.06, TEAL_DIM, f"raKT{i}")
        tb_kn = _tb(slide, kx + 0.04, 2.88, kw - 0.08, 0.28)
        _para(tb_kn.text_frame, name, 10, TEAL, bold=True, align=PP_ALIGN.CENTER)
        tb_kd = _tb(slide, kx + 0.04, 3.18, kw - 0.08, 0.36)
        _para(tb_kd.text_frame, desc, 9, MUTED, align=PP_ALIGN.CENTER)

    _rect(slide, rx, 3.76, rw, 0.56, TEAL_DARK, "raStoreNote")
    tb_sn = _tb(slide, rx + 0.12, 3.82, rw - 0.20, 0.42)
    _para(tb_sn.text_frame, "Global state: useSelector(state.chat)  ·  useDispatch() → Redux actions", 11, TEAL)

    _rect(slide, rx, 4.46, rw / 2 - 0.08, 0.56, CARD, "raBuild")
    tb_bc = _tb(slide, rx + 0.12, 4.52, rw / 2 - 0.22, 0.42)
    _para(tb_bc.text_frame, "Build: Vite 5 + SWC\nnpm run build → /dist served by FastAPI", 10, MUTED)
    _rect(slide, rx + rw / 2 + 0.04, 4.46, rw / 2 - 0.04, 0.56, CARD, "raTW")
    tb_tw = _tb(slide, rx + rw / 2 + 0.16, 4.52, rw / 2 - 0.22, 0.42)
    _para(tb_tw.text_frame, "Styling: Tailwind CSS utility classes\ndark theme, zero custom CSS files", 10, MUTED)


# ── N2: Redux Toolkit ─────────────────────────────────────────────────────────

def add_redux_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Redux Toolkit: Predictable Global State", PURPLE)

    _ic(slide, 0.5, 1.10, 5.80, 1.55, "WHAT IS REDUX TOOLKIT?", PURPLE, [
        "Single global JS object describes the entire application state",
        "Actions → Reducers → New state: unidirectional, fully predictable",
        "createSlice(): combines reducer + action creators in one file",
        "Immer built in: write mutating-looking code, get immutable state",
    ], "rd1")
    _ic(slide, 0.5, 2.77, 5.80, 1.10, "WHY NOT REACT CONTEXT?", PURPLE, [
        "Context re-renders the entire subtree on any change — too slow for chat",
        "No DevTools: Redux gives time-travel debug + action log in browser",
        "createAsyncThunk: API calls + loading/error states in one pattern",
    ], "rd2")
    _ic(slide, 0.5, 4.00, 5.80, 1.85, "DISPATCH FLOW IN OUR APP", PURPLE, [
        "User types → dispatch(appendUserMessage) → MessageList re-renders",
        "SSE event arrives → dispatch(setAnalysisProgress) → progress bar",
        "Login → dispatch(setCredentials) → all subsequent API calls are authed",
        "Token persisted in localStorage → rehydrated to authSlice on page load",
        "Settings change → dispatch(updateSettings) → next analysis uses them",
    ], "rd3")

    _divline(slide, color=PURPLE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "OUR TWO SLICES", 12, PURPLE, bold=True)

    _rect(slide, rx, 1.48, rw, 1.68, CARD, "rdAuth")
    _rect(slide, rx, 1.48, rw, 0.06, BLUE, "rdAuthT")
    tb_as = _tb(slide, rx + 0.15, 1.55, rw - 0.25, 0.32)
    _para(tb_as.text_frame, "authSlice  (features/auth/authSlice.ts)", 13, BLUE, bold=True)
    auth_fields = [
        ("user",            "{ id, username, email } | null"),
        ("token",           "JWT string | null  (persisted in localStorage)"),
        ("isAuthenticated", "boolean — gates all protected routes"),
    ]
    for i, (f, d) in enumerate(auth_fields):
        tb_f = _tb(slide, rx + 0.15, 1.93 + i * 0.30, 1.70, 0.28)
        _para(tb_f.text_frame, f, 10, BLUE, bold=True)
        tb_d = _tb(slide, rx + 1.95, 1.93 + i * 0.30, rw - 2.05, 0.28)
        _para(tb_d.text_frame, d, 10, MUTED)

    _rect(slide, rx, 3.32, rw, 2.58, CARD, "rdChat")
    _rect(slide, rx, 3.32, rw, 0.06, TEAL, "rdChatT")
    tb_cs = _tb(slide, rx + 0.15, 3.39, rw - 0.25, 0.32)
    _para(tb_cs.text_frame, "chatSlice  (features/chat/chatSlice.ts)", 13, TEAL, bold=True)
    chat_fields = [
        ("messages",        "ConversationMessage[] — full chat history"),
        ("analysisResults", "AnalysisResult | null — current result set"),
        ("isStreaming",     "boolean — spinner / SSE in progress"),
        ("streamingText",   "string — partial AI response tokens"),
        ("conversationId",  "UUID — links messages to backend"),
        ("settings",        "{ organism, cancerGenesOnly, datasetCount }"),
    ]
    for i, (f, d) in enumerate(chat_fields):
        tb_f = _tb(slide, rx + 0.15, 3.73 + i * 0.33, 1.70, 0.28)
        _para(tb_f.text_frame, f, 10, TEAL, bold=True)
        tb_d = _tb(slide, rx + 1.95, 3.73 + i * 0.33, rw - 2.05, 0.28)
        _para(tb_d.text_frame, d, 10, MUTED)

    _rect(slide, rx, 6.06, rw, 0.46, TEAL_DARK, "rdNote")
    tb_n = _tb(slide, rx + 0.12, 6.12, rw - 0.20, 0.32)
    _para(tb_n.text_frame, "Redux DevTools: replay actions, inspect state diffs, time-travel debug", 11, TEAL)


# ── N3: SSE (Frontend) ────────────────────────────────────────────────────────

def add_sse_frontend_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "SSE: It Is Just HTTP — No Special Protocol Needed", TEAL)

    # ── Left: HTTP wire format ─────────────────────────────────────────────────
    tb_lh = _tb(slide, 0.5, 1.10, 6.0, 0.30)
    _para(tb_lh.text_frame, "WHAT THE BROWSER ACTUALLY RECEIVES", 12, TEAL, bold=True)

    # Pseudo-HTTP response box
    _rect(slide, 0.5, 1.48, 6.0, 3.90, CARD, "seHTTP")
    _rect(slide, 0.5, 1.48, 6.0, 0.06, TEAL, "seHTTPT")
    tb_hdr = _tb(slide, 0.65, 1.54, 5.70, 0.28)
    _para(tb_hdr.text_frame, "HTTP RESPONSE", 10, TEAL, bold=True)

    http_lines = [
        ("HTTP/1.1 200 OK",                          LIGHT),
        ("Content-Type: text/event-stream",           BLUE),
        ("Cache-Control: no-cache",                   MUTED),
        ("Connection: keep-alive",                    MUTED),
        ("",                                          MUTED),
        ("event: searching",                          TEAL),
        ('data: {"type":"searching","gse_count":47}', AMBER),
        ("",                                          MUTED),
        ("event: analyzing",                          TEAL),
        ('data: {"gse_id":"GSE12345","progress":3}',  AMBER),
        ("",                                          MUTED),
        ("event: complete",                           TEAL),
        ('data: {"genes":[...],"n_datasets":10}',     AMBER),
    ]
    for i, (line, col) in enumerate(http_lines):
        tb_l = _tb(slide, 0.65, 1.88 + i * 0.24, 5.70, 0.24)
        _para(tb_l.text_frame, line, 10, col)

    # Key insight badge
    _rect(slide, 0.5, 5.48, 6.0, 0.56, TEAL_DARK, "seInsight")
    tb_i = _tb(slide, 0.65, 5.54, 5.70, 0.42)
    _para(tb_i.text_frame, "Persistent connection — server pushes events as they happen\nCaddy proxy: flush_interval -1  (one line, streaming enabled)", 11, TEAL)

    # SSE vs WebSocket compact comparison
    _rect(slide, 0.5, 6.14, 6.0, 0.76, CARD, "seCmp")
    for j, (label, val_sse, val_ws) in enumerate([
        ("Direction", "server → client only", "bidirectional"),
        ("Protocol",  "plain HTTP",           "ws:// upgrade"),
        ("Reconnect", "auto (browser)",        "manual"),
    ]):
        cy = 6.20 + j * 0.22
        tb_cl = _tb(slide, 0.62, cy, 1.20, 0.20)
        _para(tb_cl.text_frame, label, 8, MUTED, bold=True)
        tb_cs = _tb(slide, 1.90, cy, 2.30, 0.20)
        _para(tb_cs.text_frame, f"SSE: {val_sse}", 8, TEAL)
        tb_cw = _tb(slide, 4.30, cy, 2.10, 0.20)
        _para(tb_cw.text_frame, f"WS: {val_ws}", 8, MUTED)

    _divline(slide)
    rx, rw = 6.65, 12.633 - 6.65

    # ── Right: JavaScript client code ─────────────────────────────────────────
    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "BROWSER CLIENT  (api.ts)", 12, BLUE, bold=True)

    _rect(slide, rx, 1.48, rw, 3.10, CARD, "seJSBox")
    _rect(slide, rx, 1.48, rw, 0.06, BLUE, "seJST")
    js_lines = [
        ("// Open persistent HTTP connection", MUTED),
        ("const es = new EventSource(", LIGHT),
        ("  '/api/search/stream?q=lung+cancer'", AMBER),
        (");", LIGHT),
        ("", MUTED),
        ("es.onmessage = ({ data }) => {", LIGHT),
        ("  const event = JSON.parse(data);", LIGHT),
        ("  if (event.type === 'complete')", MUTED),
        ("    dispatch(setAnalysisResult(event));", TEAL),
        ("  else", MUTED),
        ("    dispatch(setStreamingProgress(event));", TEAL),
        ("};", LIGHT),
        ("", MUTED),
        ("// Browser handles reconnect automatically", MUTED),
        ("es.onerror = () => { /* just wait */ };", MUTED),
    ]
    for i, (line, col) in enumerate(js_lines):
        tb_j = _tb(slide, rx + 0.14, 1.56 + i * 0.18, rw - 0.20, 0.18)
        _para(tb_j.text_frame, line, 9, col)

    # Redux dispatch note
    _rect(slide, rx, 4.68, rw, 0.60, TEAL_DARK, "seDispatch")
    tb_d = _tb(slide, rx + 0.12, 4.74, rw - 0.20, 0.46)
    tf = tb_d.text_frame
    tf.word_wrap = True
    _para(tf, "Each SSE event → Redux dispatch → React re-renders automatically\nZero polling, zero timers, zero manual state sync", 10, TEAL)

    # Chat streaming comparison
    _rect(slide, rx, 5.40, rw, 0.82, CARD, "seChatNote")
    _rect(slide, rx, 5.40, rw, 0.06, CORAL, "seChatT")
    tb_cn = _tb(slide, rx + 0.14, 5.48, rw - 0.22, 0.28)
    _para(tb_cn.text_frame, "Chat uses a different mechanism:", 10, CORAL, bold=True)
    tb_cd = _tb(slide, rx + 0.14, 5.78, rw - 0.22, 0.36)
    _para(tb_cd.text_frame, "fetch() + ReadableStream → reader.read() → append tokens\nReason: chat is bidirectional (request + streamed response)", 9, MUTED)

    # Why SSE for analysis, not chat
    _rect(slide, rx, 6.36, rw, 0.54, CARD, "seReason")
    tb_re = _tb(slide, rx + 0.12, 6.42, rw - 0.20, 0.40)
    _para(tb_re.text_frame, "Analysis: one request → many events → EventSource\nChat: one request → one streamed body → ReadableStream", 9, MUTED)


# ── N4: Recharts ──────────────────────────────────────────────────────────────

def add_recharts_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Recharts: Declarative Data Visualisation for React", TEAL)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS RECHARTS?", TEAL, [
        "React-native charting library — charts are JSX components",
        "SVG-based: sharp at any scale, easy to customise with CSS",
        "Declarative API: <LineChart data={…}><Line dataKey='survival'/></LineChart>",
        "Responsive: <ResponsiveContainer width='100%'> adapts to any screen",
    ], "rc1")
    _ic(slide, 0.5, 2.72, 5.80, 1.10, "WHY NOT D3 DIRECTLY?", TEAL, [
        "D3 requires imperative DOM manipulation — fights React's virtual DOM",
        "Recharts = D3 math + React rendering: best of both worlds",
        "Tooltips, legends, axes — built-in, no extra code",
    ], "rc2")
    _ic(slide, 0.5, 3.94, 5.80, 1.90, "KEY ADVANTAGES", TEAL, [
        "Custom tooltip: <CustomTooltip> receives active payload as React props",
        "Synchronised charts: syncId='survival' links multiple KM curves",
        "Reference lines: <ReferenceLine y={0.5}> draws median survival line",
        "Click handlers: <Scatter onClick={onGeneClick}> for volcano plot genes",
        "PNG export: html2canvas on the chart ref — one button, publication ready",
    ], "rc3")

    _divline(slide)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "CHARTS WE USE", 12, TEAL, bold=True)

    charts = [
        ("KM Curve",     "LineChart\n+ ReferenceLine\n+ custom tooltip\n+ CI shading", TEAL),
        ("Volcano Plot", "ScatterChart\n+ ReferenceLine (HR=1, p=0.05)\n+ onClick → gene modal", CORAL),
        ("Forest Plot",  "ComposedChart\n+ ErrorBar (CI)\n+ Bar (effect size)\n+ diamond pooled", BLUE),
    ]
    cw = (rw - 0.30) / 3
    for i, (name, desc, col) in enumerate(charts):
        cx = rx + i * (cw + 0.15)
        _rect(slide, cx, 1.48, cw, 3.40, CARD, f"rcC{i}")
        _rect(slide, cx, 1.48, cw, 0.06, col, f"rcCT{i}")
        tb_cn = _tb(slide, cx + 0.08, 1.55, cw - 0.16, 0.36)
        _para(tb_cn.text_frame, name, 14, col, bold=True, align=PP_ALIGN.CENTER)
        tb_cd = _tb(slide, cx + 0.08, 1.96, cw - 0.16, 2.80)
        _para(tb_cd.text_frame, desc, 11, MUTED, align=PP_ALIGN.CENTER)

    _rect(slide, rx, 5.02, rw, 1.35, CARD, "rcExtra")
    tb_ex = _tb(slide, rx + 0.15, 5.10, rw - 0.25, 1.15)
    tf = tb_ex.text_frame
    tf.word_wrap = True
    _para(tf, "ADDITIONAL VISUALISATIONS", 11, TEAL, bold=True)
    _para(tf, "• BarChart — per-gene occurrence count across datasets", 10, LIGHT, new=True)
    _para(tf, "• PieChart — platform distribution breakdown", 10, LIGHT, new=True)
    _para(tf, "• All charts support CSV + PNG download via custom export toolbar", 10, LIGHT, new=True)


# ── N5: FastAPI ───────────────────────────────────────────────────────────────

def add_fastapi_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "FastAPI: 4 Problems Solved for This App", TEAL)

    problems = [
        (
            TEAL,
            "5-minute analysis — client can't wait for one HTTP response",
            'StreamingResponse(stream_gen(), media_type="text/event-stream")',
            "→ Browser EventSource reads live progress events; connection stays open the whole analysis",
        ),
        (
            AMBER,
            "Every protected route needs JWT decode — without copy-pasting code",
            "async def endpoint(user: User = Depends(get_current_user)): ...",
            "→ One dependency function, all 6 routers protected; tested once, injected everywhere",
        ),
        (
            PURPLE,
            "20 dataset downloads block a sync framework's worker thread",
            "async def search(): await orchestrator.run(query)  # loop stays free",
            "→ asyncio.gather() runs all datasets concurrently on FastAPI's native event loop",
        ),
        (
            BLUE,
            "Backend and frontend types diverge — manual sync causes silent runtime bugs",
            "class SearchRequest(BaseModel): ...  →  GET /openapi.json  →  TypeScript",
            "→ FastAPI auto-generates OpenAPI schema; TypeScript types mirror Python models",
        ),
    ]

    lx, lw = 0.35, 6.10
    card_h, card_gap = 1.28, 0.10

    for i, (color, problem, code, impact) in enumerate(problems):
        cy = 1.12 + i * (card_h + card_gap)
        _rect(slide, lx, cy, lw, card_h, CARD, f"fa{i}")
        _rect(slide, lx, cy, 0.06, card_h, color, f"faA{i}")
        tb_p = _tb(slide, lx + 0.14, cy + 0.06, lw - 0.20, 0.24)
        _para(tb_p.text_frame, problem, 9, MUTED)
        tb_c = _tb(slide, lx + 0.14, cy + 0.34, lw - 0.20, 0.30)
        _para(tb_c.text_frame, code, 10, color, bold=True, mono=True)
        tb_i = _tb(slide, lx + 0.14, cy + 0.70, lw - 0.20, 0.46)
        _para(tb_i.text_frame, impact, 9, MUTED)

    _divline(slide, color=TEAL)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "OUR ROUTER STRUCTURE", 12, TEAL, bold=True)

    routers = [
        ("/api/search/stream", "GET — SSE analysis stream\n→ GEOSurvivalWorkflowOrchestrator", TEAL),
        ("/api/results",       "GET/POST — save & list results\n→ ResultService + SQLAlchemy", BLUE),
        ("/api/chat",          "POST /conversations/{id}/messages\n→ PydanticAIService (stream)", PURPLE),
        ("/api/auth",          "POST /token · /register · /me\n→ AuthService + JWT", AMBER),
        ("/api/compare",       "POST — side-by-side analysis\n→ CompareService", CORAL),
        ("/api/enrichment",    "POST — pathway GO enrichment\n→ PathwayEnrichmentService", BLUE),
    ]
    rh, rg = 0.70, 0.10
    for i, (route, desc, col) in enumerate(routers):
        ry = 1.48 + i * (rh + rg)
        _rect(slide, rx, ry, rw, rh, CARD, f"faR{i}")
        _rect(slide, rx, ry, 0.07, rh, col, f"faRA{i}")
        tb_r = _tb(slide, rx + 0.16, ry + 0.08, 2.20, 0.30)
        _para(tb_r.text_frame, route, 11, col, bold=True, mono=True)
        tb_d = _tb(slide, rx + 0.16, ry + 0.38, rw - 0.26, 0.26)
        _para(tb_d.text_frame, desc, 10, MUTED)

    _rect(slide, rx, 6.26, rw, 0.38, TEAL_DARK, "faNote")
    tb_n = _tb(slide, rx + 0.12, 6.30, rw - 0.20, 0.28)
    _para(tb_n.text_frame, "All endpoints are async def — event loop is never blocked", 11, TEAL)


# ── N6: Pydantic ──────────────────────────────────────────────────────────────

def add_pydantic_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Pydantic: One Model, Four Jobs", BLUE)

    # Left: code block showing the model definition
    lx, lw = 0.35, 6.10

    tb_lh = _tb(slide, lx, 1.12, lw, 0.28)
    _para(tb_lh.text_frame, "DEFINE ONCE, FOUR CONSUMERS", 11, BLUE, bold=True)

    code_block = _tb(slide, lx, 1.46, lw, 2.60)
    tf = code_block.text_frame
    tf.word_wrap = False
    code_lines = [
        ("class SearchRequest(BaseModel):",          LIGHT),
        ("    query: str",                            TEAL),
        ("    organism: str = 'Homo sapiens'",        MUTED),
        ("    cancer_genes_only: bool = False",       MUTED),
        ("    dataset_count: int = Field(10,",        AMBER),
        ("                             ge=1, le=50)", AMBER),
        ("    gene_filter: list[str] = []",           MUTED),
    ]
    for k, (line, color) in enumerate(code_lines):
        _para(tf, line, 10, color, new=(k > 0), mono=True)

    jobs = [
        (CORAL,   "FastAPI validation",
                  "POST body auto-validated; dataset_count=999 → 422 automatically; zero try/except"),
        (TEAL,    "OpenAPI /docs",
                  "Full interactive schema at /docs — free, always in sync, zero extra config"),
        (BLUE,    "TypeScript mirror",
                  "openapi-typescript generates SearchRequest.ts; type errors at compile time"),
        (PURPLE,  "pydantic-ai structured output",
                  "LLM returns SurvivalEstimate(BaseModel) — no JSON parsing, no hallucinated fields"),
    ]
    job_y = 4.18
    job_h = 0.68
    job_gap = 0.10
    for i, (color, title, desc) in enumerate(jobs):
        jy = job_y + i * (job_h + job_gap)
        _rect(slide, lx, jy, lw, job_h, CARD, f"pyJ{i}")
        _rect(slide, lx, jy, 0.06, job_h, color, f"pyJA{i}")
        tb_t = _tb(slide, lx + 0.14, jy + 0.06, lw - 0.20, 0.26)
        _para(tb_t.text_frame, title, 11, color, bold=True)
        tb_d = _tb(slide, lx + 0.14, jy + 0.34, lw - 0.20, 0.28)
        _para(tb_d.text_frame, desc, 9, MUTED)

    _divline(slide, color=BLUE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "REQUEST → RESPONSE FLOW", 12, BLUE, bold=True)

    _fc(slide, [
        ("POST /api/search/stream  (raw JSON body)",  "query, organism, dataset_count, …", BLUE),
        ("FastAPI auto-validates via Pydantic",        "SearchRequest.model_validate(body) → typed obj", BLUE),
        ("Validation error  →  422 + field detail",   "e.g. dataset_count must be ≤ 50", CORAL),
        ("Valid → typed SearchRequest passed down",    "orchestrator.run(request) — fully typed arg", TEAL),
        ("Service returns list[GeneResult]",           "Pydantic serialises to JSON response automatically", TEAL),
        ("TypeScript interface matches GeneResult",    "No manual field mapping; type errors caught at build", BLUE),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="py")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "pyNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "pydantic-settings: Settings(BaseSettings) reads .env → typed config object; cached as singleton", 10, TEAL)


# ── N7: asyncio ───────────────────────────────────────────────────────────────

def add_asyncio_slide(prs: Presentation) -> None:
    """Slide 16 — asyncio production patterns. Slide 38 owns gather/semaphore/timing proof."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "asyncio: Production Patterns Beyond gather()", PURPLE)

    patterns = [
        (
            TEAL,
            "Never use requests.get() — always httpx.AsyncClient",
            "async with httpx.AsyncClient(timeout=120) as c:\n    r = await c.get(ncbi_url)",
            "requests.get() blocks the event loop — all other coroutines freeze until the download finishes",
        ),
        (
            AMBER,
            "CPU-bound code (Cox regression) must leave the event loop",
            "loop = asyncio.get_event_loop()\nresult = await loop.run_in_executor(None, fitter.fit, df)",
            "lifelines CoxPHFitter.fit() is pure CPU — blocks for 1–5 s; executor sends it to a thread pool",
        ),
        (
            PURPLE,
            "Async generators make SSE streaming a first-class citizen",
            "async def stream():\n    async for step in orchestrator.run(q):\n        yield f\"data: {json.dumps(step)}\\n\\n\"",
            "Each yield flushes immediately to the browser; no buffering; works natively with StreamingResponse",
        ),
        (
            CORAL,
            "Long loops starve other coroutines — yield at checkpoints",
            "for batch in chunks(genes, 100):\n    results.extend(process(batch))\n    await asyncio.sleep(0)  # cooperative yield",
            "asyncio is cooperative: if you never await, you own the loop; sleep(0) is the cost-free yield point",
        ),
    ]

    lx, lw = 0.35, 6.10
    card_h, card_gap = 1.34, 0.10

    for i, (color, title, code, why) in enumerate(patterns):
        cy = 1.12 + i * (card_h + card_gap)
        _rect(slide, lx, cy, lw, card_h, CARD, f"ao{i}")
        _rect(slide, lx, cy, 0.06, card_h, color, f"aoA{i}")
        tb_t = _tb(slide, lx + 0.14, cy + 0.06, lw - 0.20, 0.22)
        _para(tb_t.text_frame, title, 9, MUTED)
        tb_c = _tb(slide, lx + 0.14, cy + 0.30, lw - 0.20, 0.42)
        _para(tb_c.text_frame, code, 9, color, bold=True, mono=True)
        tb_w = _tb(slide, lx + 0.14, cy + 0.76, lw - 0.20, 0.46)
        _para(tb_w.text_frame, why, 9, MUTED)

    _divline(slide, color=PURPLE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "THE EVENT LOOP MENTAL MODEL", 12, PURPLE, bold=True)

    _ic(slide, rx, 1.48, rw, 2.20, "ONE THREAD, MANY COROUTINES", PURPLE, [
        "The event loop runs one coroutine at a time — true parallelism happens at I/O wait points",
        "await suspends the current coroutine → loop picks up another → comes back when I/O done",
        "This project: ~20 FTP downloads all suspended at await simultaneously — all progress together",
        "No OS thread switch overhead; no GIL contention; scales to hundreds of waiting coroutines",
        "Rule of thumb: if it touches the network, disk, or DB — it must be async. If it's CPU — use executor.",
    ], "aoMM")

    _ic(slide, rx, 3.80, rw, 1.45, "WHAT BREAKS ASYNC (AND HOW WE AVOID IT)", CORAL, [
        "time.sleep(n) → freezes the loop for n seconds → use await asyncio.sleep(n)",
        "requests.get() → synchronous FTP/HTTP blocks → use httpx.AsyncClient",
        "pandas/numpy heavy computation → blocks CPU → run_in_executor to thread pool",
        "Bare 'except Exception' hides broken coroutines → catch specific exceptions, log them",
    ], "aoBreak")

    _rect(slide, rx, 5.36, rw, 0.90, TEAL_DARK, "aoNote")
    tb_n = _tb(slide, rx + 0.14, 5.44, rw - 0.24, 0.72)
    _para(tb_n.text_frame,
          "asyncio.gather() timing proof + Semaphore(5) NCBI throttling on slide 38\n"
          "Key result: 20 datasets × 60 s avg  →  20 min sequential  vs  ~2 min parallel",
          11, TEAL, align=PP_ALIGN.CENTER)


# ── N8: pydantic-ai ───────────────────────────────────────────────────────────

def add_pydantic_ai_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "pydantic-ai: Production-Grade AI Agents", TEAL)

    _ic(slide, 0.5, 1.10, 5.80, 1.55, "WHAT IS pydantic-ai?", TEAL, [
        "Python framework for building LLM agents with type-safe tool calls",
        "Agent(model, tools=[…]) — declarative agent definition",
        "agent.run_stream() — async generator, yields partial response tokens",
        "Type-safe structured output: agent returns validated Pydantic model",
    ], "pa1")
    _ic(slide, 0.5, 2.77, 5.80, 1.10, "WHY NOT LANGCHAIN?", TEAL, [
        "LangChain: heavy abstractions, unpredictable internals, v1→v2 breaks",
        "pydantic-ai: minimal API surface, first-class async, transparent internals",
        "Structured output is native — no OutputParser gymnastics",
    ], "pa2")
    _ic(slide, 0.5, 4.00, 5.80, 1.85, "OUR AGENT SETUP", TEAL, [
        "Model: mistral-large-latest (default) or claude-haiku via ANTHROPIC_KEY",
        "5 tools registered via AGENT_TOOLS list in agent_tools.py",
        "Dynamic system prompt: @agent.system_prompt reads ctx.deps.user_settings",
        "History: char-count trimming ≤16 000 chars — drops oldest messages first",
        "set_deps(): called at startup to bind services; invalidates cached agents",
    ], "pa3")

    _divline(slide)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "AGENT REQUEST LIFECYCLE", 12, TEAL, bold=True)

    _fc(slide, [
        ("User sends message",              "POST /api/chat/{conv_id}/messages", BLUE),
        ("PydanticAIService.stream_message", "Builds AgentDeps(user_settings, db, …)", TEAL),
        ("agent.run_stream(content, …)",     "History + system prompt injected", PURPLE),
        ("LLM decides to call a tool",       "Tool fn executes → returns real GEO data", AMBER),
        ("LLM generates response",           "Text tokens streamed via ReadableStream", TEAL),
        ("Domain Score computed",            "_compute_domain_score() — 0–100 badge", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="pa")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "paNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Mistral embed for RAG · same pattern for Claude if ANTHROPIC_KEY is set", 10, TEAL)


# ── N9: lifelines ─────────────────────────────────────────────────────────────

def add_lifelines_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "lifelines: Survival Analysis in Python", PURPLE)

    _ic(slide, 0.5, 1.10, 5.80, 1.55, "WHAT IS lifelines?", PURPLE, [
        "Python library for survival and event-time analysis (pandas-based)",
        "KaplanMeierFitter: non-parametric survival curve from (T, E) arrays",
        "CoxPHFitter: semi-parametric hazard model — estimates HR per covariate",
        "logrank_test(): p-value for difference between two survival curves",
    ], "ll1")
    _ic(slide, 0.5, 2.77, 5.80, 1.10, "WHY lifelines?", PURPLE, [
        "scikit-learn has no survival analysis — lifelines fills the gap",
        "Pandas-native API: .fit(df, duration_col=, event_col=) — no numpy reshaping",
        "Active maintenance, clean docs, consistent with R's survival package results",
    ], "ll2")
    _ic(slide, 0.5, 3.99, 5.80, 1.85, "HOW WE USE IT", PURPLE, [
        "Median split: df['group'] = expr > expr.median() → 'high'/'low'",
        "CoxPHFitter().fit(df, 'os_time', 'os_event').params_ → HR, p, CI",
        "KaplanMeierFitter().fit(T_high, E_high) → .survival_function_",
        "logrank_test(T_high, T_low, E_high, E_low).p_value → slide 4 p-value",
        "All per-dataset: parallelised via asyncio.gather across 10–20 datasets",
    ], "ll3")

    _divline(slide, color=PURPLE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "ANALYSIS PIPELINE (per dataset)", 12, PURPLE, bold=True)

    _fc(slide, [
        ("Expression matrix (pandas DataFrame)",  "genes × samples, float values", BLUE),
        ("Probe → gene mapping applied",           "GeneMappingService → gene symbols", TEAL),
        ("Survival metadata aligned",              "os_time, os_event columns matched to samples", AMBER),
        ("Per-gene: median split → high / low",    "df[gene] > df[gene].median()", PURPLE),
        ("CoxPHFitter.fit()  →  HR, p, CI",        "~0.1 s per gene on 200-sample dataset", TEAL),
        ("KaplanMeierFitter.fit()  →  KM curve",   "Stored as time-series for frontend rendering", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="ll")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "llNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Significance filter: p < 0.05  ·  min_occurrence ≥ 2 datasets  ·  optional COSMIC gene list", 10, TEAL)


# ── N10: SQLAlchemy ───────────────────────────────────────────────────────────

def add_sqlalchemy_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "SQLAlchemy Async: One Codebase, Two Databases", BLUE)

    lx, lw = 0.35, 6.10

    tb_lh = _tb(slide, lx, 1.12, lw, 0.28)
    _para(tb_lh.text_frame, "THE ENV-VAR SWITCH", 11, BLUE, bold=True)

    # Config code block
    cfg_tb = _tb(slide, lx, 1.46, lw, 1.60)
    tf = cfg_tb.text_frame
    tf.word_wrap = False
    cfg_lines = [
        ("# config.py",                                   MUTED),
        ("DATABASE_URL = os.getenv(",                     LIGHT),
        ('    "DATABASE_URL",',                            MUTED),
        ('    "sqlite+aiosqlite:///./geo_chat.db"',        AMBER),
        (")",                                             LIGHT),
        ("# → postgresql+asyncpg://… in production",      BLUE),
        ("# Same models.  Same Alembic.  Same Depends().", TEAL),
    ]
    for k, (line, color) in enumerate(cfg_lines):
        _para(tf, line, 10, color, new=(k > 0), mono=True)

    # Switch comparison
    sw_y = 3.14
    _rect(slide, lx, sw_y, lw, 0.24, CARD_ALT, "saSwHdr")
    tb_swh = _tb(slide, lx + 0.12, sw_y + 0.04, lw - 0.20, 0.16)
    _para(tb_swh.text_frame, "DATABASE_URL not set  →  SQLite       |       DATABASE_URL=pg://…  →  PostgreSQL", 9, MUTED)

    sw_items = [
        ("Dev / CI", "sqlite+aiosqlite\n$0 · file-based\nzero infra", AMBER),
        ("Production", "postgresql+asyncpg\n$20+/mo managed\nfull MVCC", BLUE),
    ]
    sw_bw = (lw - 0.10) / 2
    for j, (label, desc, color) in enumerate(sw_items):
        sx = lx + j * (sw_bw + 0.10)
        _rect(slide, sx, sw_y + 0.28, sw_bw, 0.92, CARD, f"saSw{j}")
        _rect(slide, sx, sw_y + 0.28, sw_bw, 0.06, color, f"saSwS{j}")
        tb_sl = _tb(slide, sx + 0.10, sw_y + 0.38, sw_bw - 0.16, 0.26)
        _para(tb_sl.text_frame, label, 11, color, bold=True)
        tb_sd = _tb(slide, sx + 0.10, sw_y + 0.66, sw_bw - 0.16, 0.46)
        _para(tb_sd.text_frame, desc, 9, MUTED)

    _ic(slide, lx, 4.52, lw, 1.58, "KEY INSIGHT", TEAL, [
        "Zero code change to switch databases — only the DATABASE_URL env var changes",
        "SQLAlchemy drivers: aiosqlite (SQLite) and asyncpg (PostgreSQL) are interchangeable",
        "Alembic migrations work identically on both — uv run alembic upgrade head",
        "Switch trigger: > 50 concurrent writers, or need for PostgreSQL-specific features",
    ], "saKI")

    _divline(slide, color=BLUE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "ASYNC SESSION DEPENDENCY INJECTION", 12, BLUE, bold=True)

    # DI pattern code
    di_tb = _tb(slide, rx, 1.46, rw, 2.10)
    tf_di = di_tb.text_frame
    tf_di.word_wrap = False
    di_lines = [
        ("engine = create_async_engine(DATABASE_URL)",       BLUE),
        ("AsyncSessionLocal = async_sessionmaker(engine)",   LIGHT),
        ("",                                                  MUTED),
        ("async def get_db():",                              TEAL),
        ("    async with AsyncSessionLocal() as session:",   TEAL),
        ("        yield session  # auto-close on exit",      MUTED),
        ("",                                                  MUTED),
        ("@router.get('/results')",                          AMBER),
        ("async def list_results(",                          LIGHT),
        ("    db: AsyncSession = Depends(get_db),",          BLUE),
        ("    user: User = Depends(get_current_user),",      AMBER),
        ("):",                                               LIGHT),
    ]
    for k, (line, color) in enumerate(di_lines):
        _para(tf_di, line, 9, color, new=(k > 0), mono=True)

    _ic(slide, rx, 3.68, rw, 2.00, "OUR DB MODELS", BLUE, [
        "User — id, username, hashed_password; source of JWT sub claim",
        "AnalysisResult — id, user_id, query, result_json, created_at; history page",
        "Conversation + Message — chat history; tool_calls stored as JSON column",
        "All models use DeclarativeBase; Alembic auto-detects schema changes",
    ], "saMod")

    _rect(slide, rx, 5.80, rw, 0.44, TEAL_DARK, "saNote")
    tb_n = _tb(slide, rx + 0.12, 5.86, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "uv run alembic revision --autogenerate -m 'desc'  +  uv run alembic upgrade head", 10, TEAL, mono=True)


# ── N11: Mistral + uv ─────────────────────────────────────────────────────────

def add_mistral_uv_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Mistral Embeddings + uv: RAG Without Infrastructure", TEAL)

    divx = 6.70
    lx, lw = 0.35, divx - 0.35 - 0.20
    rx, rw = divx + 0.20, 13.133 - divx - 0.20
    _rect(slide, divx, 1.05, 0.04, 6.10, TEAL, "muDiv")

    # ── Left: Mistral Embeddings ──────────────────────────────────────────────
    tb_lh = _tb(slide, lx, 1.12, lw, 0.28)
    _para(tb_lh.text_frame, "MISTRAL EMBEDDINGS  ·  WHY THEY MATTER HERE", 11, TEAL, bold=True)

    # Concept visual: text → vector → cosine
    concept_y = 1.46
    concept_items = [
        ('"TP53 lung survival"', LIGHT),
        ("[0.23, −0.41, 0.87,\n…1024 dims…, 0.16]", TEAL),
        ("cosine similarity\n→ top-5 doc chunks", BLUE),
    ]
    cw = (lw - 2 * 0.08) / 3
    for j, (label, color) in enumerate(concept_items):
        cx = lx + j * (cw + 0.08)
        _rect(slide, cx, concept_y, cw, 0.76, CARD, f"muC{j}")
        _rect(slide, cx, concept_y, cw, 0.06, color, f"muCS{j}")
        tb_cl = _tb(slide, cx + 0.08, concept_y + 0.12, cw - 0.14, 0.52)
        _para(tb_cl.text_frame, label, 9, color, align=PP_ALIGN.CENTER, mono=(j == 1))
        if j < 2:
            _rect(slide, cx + cw + 0.01, concept_y + 0.30, 0.06, 0.12, TEAL, f"muCA{j}")

    _rect(slide, lx, 2.34, lw, 0.28, TEAL_DARK, "muKey")
    tb_key = _tb(slide, lx + 0.12, 2.38, lw - 0.20, 0.20)
    _para(tb_key.text_frame, "Anthropic has NO embedding API — Mistral fills this gap for us", 10, TEAL)

    _ic(slide, lx, 2.72, lw, 1.50, "RAG IMPLEMENTATION  (pure numpy — no pgvector, no FAISS)", TEAL, [
        "500 GEO dataset summaries embedded at startup → stored in data/rag_index.json",
        "Query time: embed question → numpy.dot(query_vec, index_matrix) → argsort → top-5",
        "1024-dim × 500 chunks = 500K floats → 2 MB in memory → lookup in ~1ms",
        "Rebuild: uv run python scripts/build_rag_index.py  (re-embeds all chunks)",
    ], "muRAG")

    _ic(slide, lx, 4.34, lw, 1.80, "WHY THIS DESIGN?", TEAL, [
        "Zero infra: no pgvector extension, no FAISS install, no vector DB service",
        "500 chunks is trivially small — numpy in-memory is faster than any DB round-trip",
        "Mistral SDK: same client used for LLM chat and embeddings — one dependency",
        "Index is transparent JSON — diff-able in git, readable without special tooling",
        "Trade-off: won't scale past ~50K chunks without switching to FAISS/pgvector",
    ], "muWhy")

    # ── Right: uv ────────────────────────────────────────────────────────────
    tb_rh = _tb(slide, rx, 1.12, rw, 0.28)
    _para(tb_rh.text_frame, "uv  ·  THE DOCKER LAYER CACHE INSIGHT", 11, AMBER, bold=True)

    docker_tb = _tb(slide, rx, 1.46, rw, 2.30)
    tf_d = docker_tb.text_frame
    tf_d.word_wrap = False
    docker_lines = [
        ("# Dockerfile (multi-stage)",              MUTED),
        ("FROM ghcr.io/astral-sh/uv AS builder",    AMBER),
        ("",                                         MUTED),
        ("COPY pyproject.toml uv.lock ./",           LIGHT),
        ("RUN uv sync --frozen --no-dev",            TEAL),
        ("# ↑ cached unless uv.lock changes",        MUTED),
        ("",                                         MUTED),
        ("COPY . .",                                 LIGHT),
        ("# ↑ app code changes don't reinstall deps", MUTED),
        ("",                                         MUTED),
        ("FROM python:3.13-slim",                    AMBER),
        ("COPY --from=builder .venv .venv",          LIGHT),
        ("COPY --from=node /dist static/",           BLUE),
    ]
    for k, (line, color) in enumerate(docker_lines):
        _para(tf_d, line, 9, color, new=(k > 0), mono=True)

    _ic(slide, rx, 3.88, rw, 1.40, "WHY uv OVER pip?", AMBER, [
        "10–100× faster install — Rust resolver vs Python resolver",
        "uv.lock locks exact versions across all environments (CI, Docker, dev)",
        "uv run scripts/… runs in the managed venv without activation ceremony",
        "In CI: cache the uv layer by hash of uv.lock — near-instant rebuilds",
    ], "muUV")

    _ic(slide, rx, 5.40, rw, 0.74, "THREE BUILD STAGES", AMBER, [
        "Node 22: npm ci + vite build → /dist  (frontend bundle)",
        "uv builder: uv sync → .venv  (Python deps, no dev extras)",
        "python:3.13-slim: copy .venv + /dist → lean 180 MB runtime image",
    ], "muStg")


# ── N12: GEOClient ────────────────────────────────────────────────────────────

def add_geoclient_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "GEOClient: Accessing 200K+ Datasets from NCBI GEO", TEAL)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS NCBI GEO?", TEAL, [
        "Gene Expression Omnibus: largest public repository of gene expression data",
        "200K+ datasets (GSE), each with expression matrix + sample metadata",
        "Access via Entrez API (eSearch/eFetch) and FTP download",
        "SOFT / series_matrix.txt.gz: compressed tabular format for each study",
    ], "gc1")
    _ic(slide, 0.5, 2.72, 5.80, 1.30, "NCBI ENTREZ API", TEAL, [
        "eSearch: text query → list of GSE IDs (up to 10K results per call)",
        "eFetch: GSE ID → GDS metadata XML (title, organism, sample count)",
        "Rate limit: 10 req/s with email header — we set EMAIL in .env",
        "NCBI SOFT format → miniML XML or tab-delimited series_matrix.txt.gz",
    ], "gc2")
    _ic(slide, 0.5, 4.14, 5.80, 1.70, "OUR GEOClient RESPONSIBILITIES", TEAL, [
        "search_datasets(query): eSearch → ranked GSE ID list",
        "fetch_dataset_info(gse_id): eFetch → title, organism, n_samples",
        "download_series_matrix(gse_id): FTP → series_matrix.txt.gz stream",
        "fetch_platform_annotation(gpl_id): FTP → GPL annotation file",
        "async with httpx.AsyncClient() — all requests non-blocking",
    ], "gc3")

    _divline(slide)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "DATA ACQUISITION FLOW", 12, TEAL, bold=True)

    _fc(slide, [
        ("query = 'lung adenocarcinoma survival'",    "User NL input (or structured params)", BLUE),
        ("eSearch → [GSE12345, GSE67890, …]",         "NCBI returns up to dataset_count GSE IDs", TEAL),
        ("GEORankingService scores each GSE",         "Filter by sample count, organism, quality", AMBER),
        ("asyncio.gather(download_series_matrix(…))", "All chosen GSEs downloaded in parallel", PURPLE),
        ("series_matrix.txt.gz → pandas DataFrame",  "Expression matrix extracted + parsed", TEAL),
        ("GPL annotation fetched (or cache hit)",     "Probe → gene symbol mapping applied", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="gc")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "gcNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Retry logic: 3 attempts with exponential backoff on NCBI 429/503 errors", 10, TEAL)


# ── N13: GEOLoaderService ─────────────────────────────────────────────────────

def add_geoloader_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "GEOLoaderService: Parsing Expression Matrices", AMBER)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "THE PARSING CHALLENGE", AMBER, [
        "series_matrix.txt.gz: custom SOFT format — not CSV, not TSV, not JSON",
        "Header: !Sample_ metadata lines interspersed with !Series_ globals",
        "Matrix block: starts with '!series_matrix_table_begin', ends with '_end'",
        "Probe IDs as row index, GSM accessions as column headers",
    ], "gl1")
    _ic(slide, 0.5, 2.72, 5.80, 1.50, "SURVIVAL METADATA DETECTION", AMBER, [
        "Sample characteristics (characteristics_ch1) contain clinical annotation",
        "Regex patterns match: overall survival, progression free survival, event",
        "pydantic-ai fallback: structured output with tool call when regex misses",
        "Hybrid approach catches ~95% of GEO datasets without expert curation",
    ], "gl2")
    _ic(slide, 0.5, 4.35, 5.80, 1.50, "EXPRESSION MATRIX PROCESSING", AMBER, [
        "pandas read_csv with sep='\\t', comment='!': skips SOFT header lines",
        "Float parsing: replace non-numeric probe annotations before conversion",
        "Multi-probe aggregation: groupby(gene_symbol).mean() after mapping",
        "Sample alignment: inner join on sample IDs between expr + survival",
    ], "gl3")

    _divline(slide, color=AMBER)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "PARSING PIPELINE", 12, AMBER, bold=True)

    _fc(slide, [
        ("series_matrix.txt.gz  (compressed SOFT)", "FTP stream → gzip.open() in memory", BLUE),
        ("Parse !Sample_characteristics_ch1",        "Extract clinical metadata per sample", AMBER),
        ("Detect os_time + os_event columns",         "Regex OR pydantic-ai structured output", TEAL),
        ("Extract table_begin…table_end block",       "pandas.read_csv(sep='\\t') → probe×sample", AMBER),
        ("GeneMappingService maps probes → genes",    "Aggregates multi-probe per gene (mean)", PURPLE),
        ("Align expr matrix + survival metadata",      "Inner join on GSM sample accessions", TEAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="gl")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "glNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Result cached as Parquet: subsequent loads from disk in ~50 ms vs 30–120 s reparse", 10, TEAL)


# ── N14: Orchestrator ─────────────────────────────────────────────────────────

def add_orchestrator_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "GEOSurvivalWorkflowOrchestrator: Pipeline Coordinator", PURPLE)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS AN ORCHESTRATOR?", PURPLE, [
        "Top-level service that sequences and coordinates other services",
        "No business logic of its own — delegates to specialised services",
        "Manages concurrency: asyncio.gather for parallel dataset processing",
        "Owns SSE event stream: yields progress events to the HTTP response",
    ], "or1")
    _ic(slide, 0.5, 2.72, 5.80, 1.50, "WHY A SEPARATE ORCHESTRATOR?", PURPLE, [
        "Separation of concerns: GEOClient doesn't know about SurvivalAnalysis",
        "Testable: mock any service, orchestrator logic stays the same",
        "In-memory result cache (OrderedDict, max 10): deduplicates re-runs",
        "Cancellation: if SSE client disconnects, asyncio task is cancelled",
    ], "or2")
    _ic(slide, 0.5, 4.35, 5.80, 1.50, "SSE STREAMING PATTERN", PURPLE, [
        "FastAPI route: async def stream() → StreamingResponse(generator())",
        "yield 'data: {\"type\": \"searching\", …}\\n\\n' — SSE format",
        "After each stage: yield progress event → browser updates immediately",
        "Final event: complete with full result JSON → dispatch to Redux store",
    ], "or3")

    _divline(slide, color=PURPLE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "ORCHESTRATION SEQUENCE", 12, PURPLE, bold=True)

    _fc(slide, [
        ("GEOClient.search_datasets(query)",       "Returns ranked GSE ID list", BLUE),
        ("GEORankingService.filter_datasets(list)", "Applies quality + organism filters", AMBER),
        ("yield event: searching  →  SSE",          "Frontend shows 'Finding datasets…'", TEAL),
        ("asyncio.gather(*[load(gse) for gse in …])", "All datasets loaded in parallel", PURPLE),
        ("yield event: analyzing  →  SSE per batch", "Frontend shows per-dataset progress", TEAL),
        ("SurvivalAnalysisService.rank_cross_cohort", "Pool HRs → return final GeneResult[]", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="or")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "orNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Result cache key: hash(query + organism + settings) — avoids repeat downloads", 10, TEAL)


# ── N15: GEORankingService ────────────────────────────────────────────────────

def add_ranking_service_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "GEORankingService: Selecting the Best Datasets", AMBER)

    _ic(slide, 0.5, 1.10, 5.80, 1.30, "THE SELECTION PROBLEM", AMBER, [
        "eSearch returns up to 200K GSE IDs for broad disease queries",
        "We can only process ~10–20 datasets in 2–5 minutes",
        "Bad dataset = low sample count, missing survival data, wrong organism",
        "Wrong selection wastes the user's time and skews the results",
    ], "rs1")
    _ic(slide, 0.5, 2.52, 5.80, 2.00, "RANKING SIGNALS (scored 0–100)", AMBER, [
        "Sample count: ≥ 50 samples = +30 pts; < 10 = eliminated",
        "Keyword match: query terms in title/abstract = +25 pts",
        "Organism: human only if organism filter set = +20 pts",
        "Platform type: known expression platforms (Affymetrix/Illumina) = +15 pts",
        "Recency: datasets from last 5 years = +10 pts",
        "Pre-computed cache: previously downloaded datasets get +5 pts bonus",
    ], "rs2")
    _ic(slide, 0.5, 4.64, 5.80, 1.20, "OUTCOME", AMBER, [
        "Top dataset_count datasets selected (default: 10, max: 50)",
        "Ranked by total score — analysis starts with best datasets first",
        "Fallback: if fewer than min_datasets qualify, threshold is relaxed",
    ], "rs3")

    _divline(slide, color=AMBER)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "SCORING BREAKDOWN", 12, AMBER, bold=True)

    signals = [
        ("Sample Count",  30, TEAL),
        ("Keyword Match", 25, BLUE),
        ("Organism",      20, PURPLE),
        ("Platform Type", 15, AMBER),
        ("Recency",       10, CORAL),
    ]
    bar_left  = rx + 2.50
    bar_max_w = rw - 2.60
    bar_y_start = 1.48
    bar_h    = 0.54
    bar_gap  = 0.16
    for i, (label, pts, col) in enumerate(signals):
        by = bar_y_start + i * (bar_h + bar_gap)
        tb_l = _tb(slide, rx, by + 0.13, 2.40, 0.30)
        _para(tb_l.text_frame, label, 11, WHITE)
        bw = (pts / 30) * bar_max_w
        _rect(slide, bar_left, by + 0.12, bw, bar_h - 0.24, col, f"rsBar{i}")
        tb_p = _tb(slide, bar_left + bw + 0.08, by + 0.13, 0.60, 0.28)
        _para(tb_p.text_frame, f"+{pts}", 11, col, bold=True)

    _rect(slide, rx, 4.58, rw, 1.70, CARD, "rsNote2")
    tb_n2 = _tb(slide, rx + 0.15, 4.66, rw - 0.25, 1.52)
    tf = tb_n2.text_frame
    tf.word_wrap = True
    _para(tf, "AFTER RANKING:", 11, AMBER, bold=True)
    _para(tf, "• Low-scoring datasets skipped entirely — never downloaded", 10, LIGHT, new=True)
    _para(tf, "• Scores logged: grep 'ranking_score' in geo_logs/app.log", 10, MUTED, new=True)
    _para(tf, "• Platform-cached datasets always preferred (disk hit = fast)", 10, MUTED, new=True)


# ── N16: SurvivalAnalysisService ──────────────────────────────────────────────

def add_survival_service_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "SurvivalAnalysisService: Statistical Core", CORAL)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "RESPONSIBILITIES", CORAL, [
        "Takes aligned expression + survival DataFrame → GeneResult list",
        "Per-dataset: CoxPH + KM fit for every gene passing quality checks",
        "Cross-cohort: aggregates per-dataset results → ranked gene list",
        "Heterogeneity: computes Cochran Q + I² for forest plot display",
    ], "sv1")
    _ic(slide, 0.5, 2.72, 5.80, 1.50, "STATISTICAL DECISIONS", CORAL, [
        "Median split: simple, reproducible, no parametric assumptions needed",
        "CoxPH assumption check: Schoenfeld residuals test (p > 0.05 required)",
        "Multiple testing: Bonferroni-adjusted p reported alongside raw p-value",
        "Min sample per group: ≥ 5 events required in each arm or gene is skipped",
    ], "sv2")
    _ic(slide, 0.5, 4.35, 5.80, 1.50, "CROSS-COHORT RANKING", CORAL, [
        "Gene score = n_significant_datasets (primary) + mean(-log10(p)) (secondary)",
        "Direction consistency check: majority of datasets must agree on HR>1 or <1",
        "Inconsistent genes flagged with high I² (≥ 75%) — shown in forest plot",
        "Output: GeneResult list sorted by score, with per-dataset HR + CI arrays",
    ], "sv3")

    _divline(slide, color=CORAL)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "INPUT → PROCESS → OUTPUT", 12, CORAL, bold=True)

    # Input box
    _rect(slide, rx, 1.48, rw, 1.00, CARD, "svIn")
    _rect(slide, rx, 1.48, rw, 0.06, BLUE, "svInT")
    tb_in = _tb(slide, rx + 0.15, 1.55, rw - 0.25, 0.30)
    _para(tb_in.text_frame, "INPUT", 12, BLUE, bold=True)
    tb_ind = _tb(slide, rx + 0.15, 1.88, rw - 0.25, 0.54)
    _para(tb_ind.text_frame, "expr_df: genes × samples  ·  surv_df: os_time + os_event per sample\ngse_id, dataset_metadata", 11, MUTED)

    _rect(slide, rx + rw / 2 - 0.04, 2.48, 0.08, 0.20, TEAL, "svA0")

    # Process box
    _rect(slide, rx, 2.68, rw, 1.50, CARD, "svProc")
    _rect(slide, rx, 2.68, rw, 0.06, CORAL, "svProcT")
    tb_pr = _tb(slide, rx + 0.15, 2.75, rw - 0.25, 0.30)
    _para(tb_pr.text_frame, "PROCESS (per gene)", 12, CORAL, bold=True)
    proc_steps = [
        "1. median split expression → high_group / low_group samples",
        "2. logrank_test(T_h, T_l, E_h, E_l) → p-value",
        "3. CoxPHFitter.fit(df[group, os_time, os_event]) → HR, 95% CI",
        "4. KaplanMeierFitter.fit(T_h, E_h) → km_high; same for low",
    ]
    tb_ps = _tb(slide, rx + 0.15, 3.08, rw - 0.25, 1.04)
    tf = tb_ps.text_frame
    tf.word_wrap = True
    for i, s in enumerate(proc_steps):
        _para(tf, s, 10, LIGHT, new=(i > 0))

    _rect(slide, rx + rw / 2 - 0.04, 4.18, 0.08, 0.20, TEAL, "svA1")

    # Output box
    _rect(slide, rx, 4.38, rw, 1.02, CARD, "svOut")
    _rect(slide, rx, 4.38, rw, 0.06, TEAL, "svOutT")
    tb_ot = _tb(slide, rx + 0.15, 4.45, rw - 0.25, 0.30)
    _para(tb_ot.text_frame, "OUTPUT (per dataset)", 12, TEAL, bold=True)
    tb_od = _tb(slide, rx + 0.15, 4.78, rw - 0.25, 0.56)
    _para(tb_od.text_frame, "SurvivalResult: { gene, hr, ci_low, ci_high, pval, km_high, km_low, n_samples }", 10, MUTED)

    _rect(slide, rx, 5.52, rw, 0.44, TEAL_DARK, "svAgg")
    tb_ag = _tb(slide, rx + 0.12, 5.58, rw - 0.20, 0.30)
    _para(tb_ag.text_frame, "Cross-cohort aggregation: GeneResult[].n_significant sorted desc → ranked gene table", 10, TEAL)


# ── N17: Chat Tools ───────────────────────────────────────────────────────────

def add_chat_tools_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "AI Chat: 5 Tools Grounding Every Response in Real Data", TEAL)

    tools = [
        ("search_known_datasets",
         "WHAT: searches cached analysis history for matching GEO datasets\nWHY: instant answers from already-processed data — no re-download\nRETURNS: { gse_id, title, organism, n_genes_significant }",
         TEAL),
        ("search_geo_datasets",
         "WHAT: queries NCBI Entrez eSearch live for new GSE IDs\nWHY: when user asks about a disease not yet in cache\nRETURNS: [ { gse_id, title, n_samples, platform } ] up to 10",
         BLUE),
        ("get_gene_info",
         "WHAT: looks up gene in analysis results for HR, p-value, dataset count\nWHY: grounds gene biology questions in actual survival data\nRETURNS: { gene, hr, pval, n_datasets, km_data, direction }",
         PURPLE),
        ("estimate_query",
         "WHAT: predicts dataset count + analysis time before running\nWHY: sets expectations — 'I found 12 datasets, ~3 min'\nRETURNS: { n_datasets_est, time_est_seconds, confidence }",
         AMBER),
        ("get_user_recent_results",
         "WHAT: retrieves last 5 analyses from the user's history\nWHY: enables follow-up questions referencing a previous run\nRETURNS: [ { query, result_id, top_genes, created_at } ]",
         CORAL),
    ]

    th = (6.20) / len(tools)
    for i, (name, desc, col) in enumerate(tools):
        ty = 1.10 + i * (th + 0.05)
        _rect(slide, 0.5, ty, 12.333, th, CARD, f"ct{i}")
        _rect(slide, 0.5, ty, 0.07, th, col, f"ctA{i}")
        tb_n = _tb(slide, 0.66, ty + 0.06, 3.00, 0.32)
        _para(tb_n.text_frame, name, 12, col, bold=True)
        tb_d = _tb(slide, 0.66, ty + 0.40, 11.90, th - 0.48)
        _para(tb_d.text_frame, desc, 10, MUTED)

    _rect(slide, 0.5, 7.00, 12.333, 0.34, TEAL_DARK, "ctRule")
    tb_r = _tb(slide, 0.65, 7.04, 12.0, 0.24)
    _para(tb_r.text_frame, "Rule: every substantive user question must call ≥ 1 tool — responses without tool calls have DS = 0  (domain score badge is red)", 10, TEAL)


# ── N18: RAG ──────────────────────────────────────────────────────────────────

def add_rag_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "RAG: Retrieval-Augmented Generation with numpy", BLUE)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS RAG?", BLUE, [
        "Retrieval-Augmented Generation: add relevant context to LLM prompts",
        "Avoids hallucination by grounding answers in real retrieved documents",
        "Two phases: indexing (at startup) + retrieval (at query time)",
        "Our docs: GEO dataset summaries + gene annotation chunks (~300 entries)",
    ], "rg1")
    _ic(slide, 0.5, 2.72, 5.80, 1.10, "WHY NOT PGVECTOR / FAISS?", BLUE, [
        "Our corpus is small (~300 chunks) — numpy is 100% sufficient",
        "Zero infrastructure: no Postgres extension, no C library install",
        "numpy.dot(query_vec, doc_matrix.T) → cosine similarities in one line",
    ], "rg2")
    _ic(slide, 0.5, 3.94, 5.80, 1.90, "INDEXING PHASE (at startup)", BLUE, [
        "Load rag_docs.json: list of { id, content, metadata } documents",
        "Call Mistral embed API: batch encode all content strings → 1024-dim vectors",
        "Save to rag_index.json: { vectors: [[…], …], docs: [{…}, …] }",
        "Loaded into memory as numpy array on server startup — never re-computed",
        "Rebuild: uv run python scripts/build_rag_index.py (run after adding docs)",
    ], "rg3")

    _divline(slide, color=BLUE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "RETRIEVAL PHASE (per chat message)", 12, BLUE, bold=True)

    _fc(slide, [
        ("User sends message",               "chat: 'What is TP53 HR in lung cancer?'", BLUE),
        ("Embed query with mistral-embed",    "1024-dim vector via Mistral API", TEAL),
        ("numpy.dot(q, index_matrix.T)",      "Cosine similarity vs all ~300 docs", PURPLE),
        ("Top-5 chunks by similarity score",  "GEO summaries most relevant to query", AMBER),
        ("Inject into system prompt",         "'Context: [doc1 text] [doc2 text] …'", TEAL),
        ("LLM generates grounded response",   "Cites real GSE IDs from retrieved context", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="rg")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "rgNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Index location: backend/data/rag_index.json  ·  Doc sources: backend/data/rag_docs.json", 10, TEAL)


# ── N19: Domain Score ─────────────────────────────────────────────────────────

def add_domain_score_slide(prs: Presentation) -> None:
    """Side-by-side contrast: generic AI chat vs this app, with scoring formula."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Domain Score: This App vs Generic AI Chat", AMBER)

    panel_y = 1.12
    panel_h = 3.80
    lx, lw = 0.35, 6.10
    rx, rw = 6.70, 6.28

    # ── Left panel: Generic AI ──────────────────────────────────────────────
    _rect(slide, lx, panel_y, lw, panel_h, CARD, "dsLP")
    _rect(slide, lx, panel_y, lw, 0.06, CORAL, "dsLTop")
    tb_lh = _tb(slide, lx + 0.14, panel_y + 0.10, lw - 0.24, 0.28)
    _para(tb_lh.text_frame, "Generic AI Chat  (e.g. ChatGPT.com)", 12, CORAL, bold=True)

    _rect(slide, lx + 0.12, panel_y + 0.46, lw - 0.24, 0.28, CARD_ALT, "dsLQ")
    tb_lq = _tb(slide, lx + 0.22, panel_y + 0.52, lw - 0.40, 0.18)
    _para(tb_lq.text_frame, "User: Does TP53 affect lung cancer survival?", 9, MUTED)

    _rect(slide, lx + 0.12, panel_y + 0.82, lw - 0.24, 1.20, CARD_ALT, "dsLA")
    tb_la = _tb(slide, lx + 0.22, panel_y + 0.92, lw - 0.40, 1.00)
    tf_la = tb_la.text_frame
    tf_la.word_wrap = True
    _para(tf_la,
          '"TP53 is a tumor suppressor gene involved in cell cycle regulation '
          'and apoptosis. Studies suggest it may influence lung cancer prognosis, '
          'though results vary across different patient populations..."', 9, MUTED)

    for i, text in enumerate(["Tools called: 0", "Dataset citations: none", "Statistical values: none"]):
        sy = panel_y + 2.14 + i * 0.28
        _rect(slide, lx + 0.12, sy, lw - 0.24, 0.24, CORAL_DIM, f"dsLS{i}")
        tb_si = _tb(slide, lx + 0.22, sy + 0.04, lw - 0.44, 0.18)
        _para(tb_si.text_frame, text, 9, CORAL)

    _rect(slide, lx + 0.12, panel_y + 3.04, lw - 0.24, 0.56, CORAL_DIM, "dsLBadge")
    tb_lb = _tb(slide, lx + 0.12, panel_y + 3.10, lw - 0.24, 0.42)
    _para(tb_lb.text_frame, "DOMAIN SCORE:  5 / 100", 18, CORAL, bold=True, align=PP_ALIGN.CENTER)
    _para(tb_lb.text_frame, "could be produced by any chatbot — zero product value", 9, MUTED,
          new=True, align=PP_ALIGN.CENTER)

    # ── Right panel: GEO app ────────────────────────────────────────────────
    _rect(slide, rx, panel_y, rw, panel_h, CARD, "dsRP")
    _rect(slide, rx, panel_y, rw, 0.06, TEAL, "dsRTop")
    tb_rh = _tb(slide, rx + 0.14, panel_y + 0.10, rw - 0.24, 0.28)
    _para(tb_rh.text_frame, "GEO Survival Analysis Chat", 12, TEAL, bold=True)

    _rect(slide, rx + 0.12, panel_y + 0.46, rw - 0.24, 0.28, CARD_ALT, "dsRQ")
    tb_rq = _tb(slide, rx + 0.22, panel_y + 0.52, rw - 0.40, 0.18)
    _para(tb_rq.text_frame, "User: Does TP53 affect lung cancer survival?", 9, MUTED)

    _rect(slide, rx + 0.12, panel_y + 0.82, rw - 0.24, 1.20, CARD_ALT, "dsRA")
    tb_ra = _tb(slide, rx + 0.22, panel_y + 0.92, rw - 0.40, 1.00)
    tf_ra = tb_ra.text_frame
    tf_ra.word_wrap = True
    _para(tf_ra,
          '"I searched 8 GEO datasets (GSE12345, GSE67890 + 6 more). In LUAD, '
          'TP53 high expression → HR=2.3 [CI: 1.4–3.2] (p=0.001, n=87 avg). '
          'Effect is consistent across cohorts (I²=21%)."', 9, LIGHT)

    right_signals = [
        ("Tools called: 2  → +40 pts",                      TEAL),
        ("GSE IDs: GSE12345, GSE67890  → +30 pts",          BLUE),
        ("HR=2.3, p=0.001  (+15)  ·  TP53 match  (+15)",    AMBER),
    ]
    for i, (text, color) in enumerate(right_signals):
        sy = panel_y + 2.14 + i * 0.28
        _rect(slide, rx + 0.12, sy, rw - 0.24, 0.24, TEAL_DARK, f"dsRS{i}")
        tb_si = _tb(slide, rx + 0.22, sy + 0.04, rw - 0.44, 0.18)
        _para(tb_si.text_frame, text, 9, color)

    _rect(slide, rx + 0.12, panel_y + 3.04, rw - 0.24, 0.56, TEAL_DARK, "dsRBadge")
    tb_rb = _tb(slide, rx + 0.12, panel_y + 3.10, rw - 0.24, 0.42)
    _para(tb_rb.text_frame, "DOMAIN SCORE:  80 / 100", 18, TEAL, bold=True, align=PP_ALIGN.CENTER)
    _para(tb_rb.text_frame, "grounded in real GEO patient data — unique product value", 9, MUTED,
          new=True, align=PP_ALIGN.CENTER)

    # ── Scoring formula strip ───────────────────────────────────────────────
    strip_y = panel_y + panel_h + 0.14
    strip_h = 7.10 - strip_y

    tb_sth = _tb(slide, 0.35, strip_y + 0.04, 9.0, 0.26)
    _para(tb_sth.text_frame,
          "SCORING FORMULA  (zero-cost: pure Python string analysis, no API call)",
          10, AMBER, bold=True)

    scoring = [
        ("Each tool called",        "+20 pts  (max 40)", TEAL),
        ("Each GSE ID cited",       "+15 pts  (max 30)", BLUE),
        ("HR / p-value / n= in text", "+15 pts",         PURPLE),
        ("Organism or gene matched", "+15 pts",           AMBER),
    ]
    n = len(scoring)
    sw = (13.133 - 0.70 - (n - 1) * 0.10) / n
    for j, (signal, pts, color) in enumerate(scoring):
        sx = 0.35 + j * (sw + 0.10)
        _rect(slide, sx, strip_y + 0.36, sw, strip_h - 0.40, CARD, f"dsSig{j}")
        _rect(slide, sx, strip_y + 0.36, sw, 0.06, color, f"dsSigS{j}")
        tb_ss = _tb(slide, sx + 0.12, strip_y + 0.48, sw - 0.22, 0.26)
        _para(tb_ss.text_frame, signal, 10, color, bold=True)
        tb_sp = _tb(slide, sx + 0.12, strip_y + 0.80, sw - 0.22, 0.30)
        _para(tb_sp.text_frame, pts, 12, color, bold=True)


# ── N20: SSE Pipeline (end-to-end) ────────────────────────────────────────────

def add_sse_pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "End-to-End SSE Pipeline: From Query to Live Results", TEAL)

    stages = [
        ("Browser",         "User submits NL query\nnew EventSource('/api/search/stream?q=…')", BLUE),
        ("Caddy Proxy",     "flush_interval -1\nenables streaming pass-through", TEAL),
        ("FastAPI Router",  "/api/search/stream\nStreamingResponse(generator())", TEAL),
        ("Orchestrator",    "asyncio.gather downloads\nyield SSE events per stage", PURPLE),
        ("Services",        "GEOClient + GeneMapping\n+ SurvivalAnalysis", AMBER),
        ("Redux",           "dispatch(setProgress)\ncomponent re-renders", CORAL),
    ]
    total_w = 12.333
    bw = (total_w - 0.50 - 0.30) / len(stages)
    gap = 0.06
    for i, (name, desc, col) in enumerate(stages):
        bx = 0.5 + i * (bw + gap)
        _rect(slide, bx, 1.10, bw, 1.55, CARD, f"sp{i}")
        _rect(slide, bx, 1.10, bw, 0.06, col, f"spT{i}")
        tb_n = _tb(slide, bx + 0.05, 1.18, bw - 0.10, 0.32)
        _para(tb_n.text_frame, name, 12, col, bold=True, align=PP_ALIGN.CENTER)
        tb_d = _tb(slide, bx + 0.05, 1.54, bw - 0.10, 1.04)
        _para(tb_d.text_frame, desc, 9, MUTED, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            ax = bx + bw + 0.01
            _rect(slide, ax, 1.10 + 1.55 / 2 - 0.03, gap - 0.02, 0.06, TEAL, f"spA{i}")

    # Event stream
    tb_evh = _tb(slide, 0.5, 2.82, 12.333, 0.30)
    _para(tb_evh.text_frame, "SSE EVENT STREAM  (browser sees these in sequence)", 12, TEAL, bold=True)

    events = [
        ("searching",    "Finding datasets…", BLUE),
        ("ranking",      "Scoring 47 GSEs…", TEAL),
        ("downloading",  "GET GSE12345.gz", PURPLE),
        ("analyzing",    "Cox + KM running…", AMBER),
        ("progress",     "Gene 1234 / 18000", CORAL),
        ("complete",     "{ genes: […] }", TEAL),
    ]
    ew = (total_w - 0.50) / len(events) - 0.05
    for i, (etype, data, col) in enumerate(events):
        ex = 0.5 + i * (ew + 0.05)
        _rect(slide, ex, 3.20, ew, 1.10, CARD, f"spEv{i}")
        _rect(slide, ex, 3.20, ew, 0.06, col, f"spEvT{i}")
        tb_et = _tb(slide, ex + 0.05, 3.27, ew - 0.10, 0.28)
        _para(tb_et.text_frame, f"event: {etype}", 9, col, bold=True, align=PP_ALIGN.CENTER)
        tb_ed = _tb(slide, ex + 0.05, 3.58, ew - 0.10, 0.64)
        _para(tb_ed.text_frame, f"data: {data}", 9, MUTED, align=PP_ALIGN.CENTER)

    # Frontend handler
    _rect(slide, 0.5, 4.46, 12.333, 1.60, CARD, "spHandler")
    _rect(slide, 0.5, 4.46, 12.333, 0.06, BLUE, "spHandlerT")
    tb_fh = _tb(slide, 0.65, 4.52, 4.0, 0.30)
    _para(tb_fh.text_frame, "FRONTEND EVENT HANDLER  (api.ts)", 11, BLUE, bold=True)
    tb_code = _tb(slide, 0.65, 4.86, 11.80, 1.12)
    _para(tb_code.text_frame,
          "es.onmessage = (e) => {\n  const { type, data } = JSON.parse(e.data);\n  if (type === 'complete') dispatch(setAnalysisResult(data));\n  else dispatch(setStreamingProgress({ type, data }));\n};",
          10, AMBER)

    _rect(slide, 0.5, 6.22, 12.333, 0.44, TEAL_DARK, "spNote")
    tb_n = _tb(slide, 0.65, 6.28, 12.0, 0.30)
    _para(tb_n.text_frame, "Reconnect: EventSource auto-retries on connection drop — user never needs to refresh", 11, TEAL)


# ── N21: JWT Auth ─────────────────────────────────────────────────────────────

def add_jwt_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "JWT Authentication: Stateless Security", BLUE)

    _ic(slide, 0.5, 1.10, 5.80, 1.50, "WHAT IS JWT?", BLUE, [
        "JSON Web Token: signed string encoding user claims (sub, exp, role)",
        "Stateless: server verifies signature — no session store needed",
        "Structure: base64(header) + '.' + base64(payload) + '.' + signature",
        "We sign with HS256 + JWT_SECRET_KEY from .env (min 32 chars)",
    ], "jw1")
    _ic(slide, 0.5, 2.72, 5.80, 1.10, "WHY JWT FOR THIS PROJECT?", BLUE, [
        "Single-server deployment: no Redis/session DB needed → simpler infra",
        "Shareable permalinks are public (/results/:id) — JWT gates write ops",
        "Token expiry: 30-day default, configurable via ACCESS_TOKEN_EXPIRE_DAYS",
    ], "jw2")
    _ic(slide, 0.5, 3.94, 5.80, 1.90, "SECURITY DECISIONS", BLUE, [
        "Passwords hashed with Argon2id via pwdlib — not bcrypt (stronger)",
        "Token stored in localStorage: XSS risk mitigated by strict CSP header",
        "Authorization: Bearer {token} header on all protected API calls",
        "get_current_user() Depends(): decodes + validates token per request",
        "Admin routes require role='admin' claim in token payload",
    ], "jw3")

    _divline(slide, color=BLUE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "AUTH FLOW", 12, BLUE, bold=True)

    _fc(slide, [
        ("POST /api/auth/register",     "username + password → Argon2id hash stored", TEAL),
        ("POST /api/auth/token",        "username + password verified → JWT issued (30d)", BLUE),
        ("Frontend stores token",       "localStorage.setItem('token', jwt)", AMBER),
        ("Subsequent API requests",      "Authorization: Bearer {token} header", PURPLE),
        ("get_current_user() Depends()", "jwt.decode(token, SECRET, ['HS256']) → User", TEAL),
        ("Expired/invalid token",        "HTTPException(401) → frontend redirects to login", CORAL),
    ], rx, 1.48, rw, 0.68, gap=0.10, pfx="jw")

    _rect(slide, rx, 5.96, rw, 0.44, TEAL_DARK, "jwNote")
    tb_n = _tb(slide, rx + 0.12, 6.02, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "Public routes: GET /api/results/{id} (shared permalink) — no token required", 10, TEAL)


# ── N22: Data Pipeline ────────────────────────────────────────────────────────

def add_data_pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Full Data Pipeline: NL Query → Publication Export", TEAL)

    stages = [
        ("NL Query",      "'lung adenocarcinoma\noverall survival'",                "User Input",    BLUE),
        ("NCBI eSearch",  "GSE IDs ranked\nby quality signals",                     "GEOClient",     TEAL),
        ("Download",      "series_matrix.txt.gz\nasyncio.gather()",                 "GEOClient",     PURPLE),
        ("Parse",         "SOFT → DataFrame\nsurvival metadata",                    "GEOLoader",     AMBER),
        ("Map Probes",    "GPL annotation\nprobe → gene symbol",                    "GeneMapping",   CORAL),
        ("Statistics",    "Cox PH + KM\nper gene per dataset",                      "SurvivalSvc",   TEAL),
        ("Cross-Cohort",  "Rank by n_sig\npool HRs + I²",                           "SurvivalSvc",   BLUE),
        ("Export",        "Ranked genes\nKM curves\nForest plots\nCSV + ZIP",       "Frontend",      TEAL),
    ]

    n = len(stages)
    bw = (12.333 - 0.50) / n - 0.08
    for i, (name, desc, svc, col) in enumerate(stages):
        bx = 0.5 + i * (bw + 0.08)
        _rect(slide, bx, 1.10, bw, 2.20, CARD, f"dp{i}")
        _rect(slide, bx, 1.10, bw, 0.06, col, f"dpT{i}")
        tb_n = _tb(slide, bx + 0.05, 1.17, bw - 0.10, 0.32)
        _para(tb_n.text_frame, name, 10, col, bold=True, align=PP_ALIGN.CENTER)
        tb_d = _tb(slide, bx + 0.05, 1.52, bw - 0.10, 1.10)
        _para(tb_d.text_frame, desc, 9, MUTED, align=PP_ALIGN.CENTER)
        tb_s = _tb(slide, bx + 0.05, 2.76, bw - 0.10, 0.28)
        _para(tb_s.text_frame, svc, 8, col, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = bx + bw + 0.01
            _rect(slide, ax, 1.10 + 2.20 / 2 - 0.03, 0.07, 0.06, TEAL, f"dpA{i}")

    # Data formats row
    tb_fmth = _tb(slide, 0.5, 3.22, 12.333, 0.30)
    _para(tb_fmth.text_frame, "DATA FORMATS AT EACH STAGE", 12, MUTED, bold=True)

    formats = [
        ("string",           BLUE),
        ("List[str]",        TEAL),
        (".gz stream",       PURPLE),
        ("DataFrame",        AMBER),
        ("DataFrame",        CORAL),
        ("SurvivalResult[]", TEAL),
        ("GeneResult[]",     BLUE),
        ("JSON + CSV",       TEAL),
    ]
    for i, (fmt, col) in enumerate(formats):
        fx = 0.5 + i * (bw + 0.08)
        _rect(slide, fx, 3.56, bw, 0.46, CARD, f"dpF{i}")
        tb_f = _tb(slide, fx + 0.04, 3.62, bw - 0.08, 0.32)
        _para(tb_f.text_frame, fmt, 9, col, align=PP_ALIGN.CENTER)

    # Cache hits
    _rect(slide, 0.5, 4.18, 12.333, 0.76, TEAL_DARK, "dpCache")
    tb_c = _tb(slide, 0.65, 4.26, 12.0, 0.60)
    tf = tb_c.text_frame
    tf.word_wrap = True
    _para(tf, "CACHE HITS (skip stage entirely):", 11, TEAL, bold=True)
    _para(tf, "Stage 3 (download) → Parquet hit: 50 ms  ·  Stage 5 (probe mapping) → LRU hit: 0 ms  ·  Stage 6 (stats) → result cache: instant", 10, LIGHT, new=True)

    # Timing
    _rect(slide, 0.5, 5.08, 12.333, 0.62, CARD, "dpTime")
    timings = [("Cold run", "2–5 min"), ("Warm (datasets cached)", "20–60 s"), ("Hot (result cached)", "< 1 s")]
    for i, (label, t) in enumerate(timings):
        tx = 0.5 + i * 4.111
        tb_t = _tb(slide, tx + 0.15, 5.14, 3.80, 0.26)
        _para(tb_t.text_frame, f"{label}: {t}", 12, TEAL if i == 2 else (BLUE if i == 0 else AMBER), bold=True)


# ── N23: Caching Layers ───────────────────────────────────────────────────────

def add_caching_layers_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Multi-Layer Caching: Architecture Overview", TEAL)

    tb_sub = _tb(slide, 0.5, 1.02, 12.333, 0.28)
    _para(tb_sub.text_frame, "Six independent cache layers — each prevents a different expensive operation", 13, MUTED)

    layers = [
        ("Result Cache",     "OrderedDict max 10",  "GEOSurvivalWorkflowOrchestrator",
         "Skip full re-run for same query/settings", "in-memory, keyed by hash(query+settings)", "restart clears it",    TEAL),
        ("Gene Mapping LRU", "OrderedDict max 3",   "GeneMappingService",
         "Skip platform re-parse on hot platforms",  "in-memory, platform_id → DataFrame",       "restart clears it",    BLUE),
        ("Gene Mapping Disk","Parquet files",        "backend/platform_mappings/",
         "Skip NCBI FTP re-download (2–30 min)",     "one .parquet per GPL ID",                  "persists forever",     PURPLE),
        ("Dataset Cache",    "Parquet files",        "backend/datasets/",
         "Skip series_matrix re-parse (30–120 s)",   "one .parquet per GSE ID",                  "persists forever",     AMBER),
        ("RAG Index",        "numpy array",          "backend/data/rag_index.json",
         "Avoid re-embedding 300 docs on each start","loaded to RAM at startup",                  "rebuild on doc change",CORAL),
        ("Conversation Fetch","TTL 30 s",            "chatSlice.ts (frontend)",
         "Skip repeated GET /conversations API calls","timestamp gated in Redux Toolkit",          "cleared on new message",BLUE),
    ]

    lh = (6.20 - 0.20) / len(layers)
    col_widths = [1.65, 1.20, 2.10, 2.60, 2.00, 1.45]
    col_headers = ["Layer", "Bound", "Location", "Prevents", "Key", "Lifetime"]
    col_x = [0.5, 2.20, 3.45, 5.60, 8.25, 10.30]

    # Header row
    for ci, (hdr, cx, cw) in enumerate(zip(col_headers, col_x, col_widths)):
        _rect(slide, cx, 1.38, cw, 0.38, TEAL_DARK, f"caHdr{ci}")
        tb_h = _tb(slide, cx + 0.06, 1.42, cw - 0.12, 0.28)
        _para(tb_h.text_frame, hdr, 10, TEAL, bold=True)

    for ri, (name, bound, loc, prevents, key, lifetime, col) in enumerate(layers):
        ry = 1.82 + ri * (lh + 0.04)
        _rect(slide, 0.5, ry, 11.78, lh, CARD if ri % 2 == 0 else CARD_ALT, f"caRow{ri}")
        vals = [name, bound, loc, prevents, key, lifetime]
        for ci, (val, cx, cw) in enumerate(zip(vals, col_x, col_widths)):
            tb_v = _tb(slide, cx + 0.06, ry + 0.06, cw - 0.12, lh - 0.10)
            text_col = col if ci == 0 else (MUTED if ci > 1 else LIGHT)
            bold = ci == 0
            _para(tb_v.text_frame, val, 9, text_col, bold=bold)


# ── N24: asyncio.gather Deep Dive ─────────────────────────────────────────────

def add_asyncio_gather_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "asyncio.gather: The 10× Speedup in Practice", PURPLE)

    _ic(slide, 0.5, 1.10, 5.80, 1.55, "WHAT asyncio.gather() DOES", PURPLE, [
        "Schedules N coroutines concurrently on the same event loop thread",
        "Returns when ALL coroutines complete — results in submission order",
        "If any raises, gather raises immediately (fail-fast by default)",
        "return_exceptions=True: collect errors, continue remaining coroutines",
    ], "ag1")
    _ic(slide, 0.5, 2.77, 5.80, 1.55, "OUR USAGE PATTERN", PURPLE, [
        "tasks = [load_dataset(gse_id) for gse_id in selected_gse_ids]",
        "results = await asyncio.gather(*tasks, return_exceptions=True)",
        "Filter: valid = [r for r in results if not isinstance(r, Exception)]",
        "Semaphore: asyncio.Semaphore(5) caps concurrent FTP downloads",
    ], "ag2")
    _ic(slide, 0.5, 4.44, 5.80, 1.40, "WHY NOT THREADS?", PURPLE, [
        "GIL: Python threads don't parallelise CPU work anyway",
        "asyncio: no GIL issue for I/O tasks — all waiting is non-blocking",
        "Memory: one OS thread vs N OS threads — lower overhead at scale",
        "Debugging: deterministic execution order vs race conditions with threads",
    ], "ag3")

    _divline(slide, color=PURPLE)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "TIMING COMPARISON (20 datasets, avg 60 s each)", 12, PURPLE, bold=True)

    # Sequential
    _rect(slide, rx, 1.48, rw, 0.26, CARD_ALT, "agSH")
    tb_sq = _tb(slide, rx + 0.10, 1.52, rw - 0.20, 0.20)
    _para(tb_sq.text_frame, "SEQUENTIAL  (20 min total)", 11, CORAL, bold=True)

    for i in range(10):
        bx = rx + i * (rw / 10)
        _rect(slide, bx, 1.78, rw / 10 - 0.05, 0.52, CORAL, f"agSB{i}")
    tb_sc = _tb(slide, rx + rw - 1.20, 1.82, 1.10, 0.30)
    _para(tb_sc.text_frame, "→ 20 min", 11, CORAL, bold=True)

    # Parallel
    _rect(slide, rx, 2.46, rw, 0.26, CARD_ALT, "agPH")
    tb_pa = _tb(slide, rx + 0.10, 2.50, rw - 0.20, 0.20)
    _para(tb_pa.text_frame, "PARALLEL  asyncio.gather()  (2 min total)", 11, TEAL, bold=True)

    for i in range(10):
        bx = rx + i * 0.22
        _rect(slide, bx, 2.76, 0.18, 1.54, TEAL, f"agPB{i}")
    _rect(slide, rx + 10 * 0.22 + 0.05, 2.76, rw - 10 * 0.22 - 0.10, 1.54, TEAL_DIM, "agPBg")
    tb_pc = _tb(slide, rx + 10 * 0.22 + 0.20, 3.20, rw - 10 * 0.22 - 0.25, 0.60)
    _para(tb_pc.text_frame, "All 20 running\nconcurrently", 12, TEAL, bold=True, align=PP_ALIGN.CENTER)
    tb_pt = _tb(slide, rx + rw - 1.20, 2.80, 1.10, 0.30)
    _para(tb_pt.text_frame, "→ 2 min", 11, TEAL, bold=True)

    # Semaphore note
    _rect(slide, rx, 4.48, rw, 0.72, CARD, "agSem")
    tb_sem = _tb(slide, rx + 0.12, 4.54, rw - 0.22, 0.56)
    tf = tb_sem.text_frame
    tf.word_wrap = True
    _para(tf, "Semaphore(5): throttles to max 5 concurrent FTP connections — avoids NCBI rate-limit ban", 11, AMBER)

    _rect(slide, rx, 5.34, rw, 0.58, TEAL_DARK, "agWin")
    tb_w = _tb(slide, rx + 0.12, 5.40, rw - 0.22, 0.42)
    _para(tb_w.text_frame, "10× wall-time reduction: the single biggest engineering win in this project", 13, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── N25: Error Handling ───────────────────────────────────────────────────────

def add_error_handling_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _title_block(slide, "Error Handling: Resilient by Design", CORAL)

    _ic(slide, 0.5, 1.10, 5.80, 1.55, "PHILOSOPHY", CORAL, [
        "Catch specific exceptions — never bare except: or except Exception:",
        "Network errors are expected: NCBI FTP is unreliable, retry 3× with backoff",
        "Graceful degradation: one dataset fails → skip it, continue with others",
        "SSE errors: yield error event to browser — user sees what went wrong",
    ], "eh1")
    _ic(slide, 0.5, 2.77, 5.80, 1.55, "BOUNDARY VALIDATION ONLY", CORAL, [
        "Trust internal code — only validate at system boundaries (user input, API)",
        "FastAPI + Pydantic: invalid request → 422 with field-level error detail",
        "JWT auth: missing/expired token → 401 HTTPException (not AssertionError)",
        "No defensive coding inside service layer — trust typed inputs",
    ], "eh2")
    _ic(slide, 0.5, 4.44, 5.80, 1.40, "LOGGING STRATEGY", CORAL, [
        "logger.info() for normal flow events (analysis start, cache hit)",
        "logger.warning() for recoverable errors (dataset skip, retry attempt)",
        "logger.exception() for unexpected errors (unhandled exception in handler)",
        "Structured format: key=value pairs for grep-able log analysis",
    ], "eh3")

    _divline(slide, color=CORAL)
    rx, rw = 6.55, 12.633 - 6.55

    tb_rh = _tb(slide, rx, 1.10, rw, 0.30)
    _para(tb_rh.text_frame, "ERROR PATHS + RESPONSES", 12, CORAL, bold=True)

    errors = [
        ("NCBI FTP timeout",       "httpx.TimeoutException",    "Retry 3× with exponential backoff (2/4/8 s)",   AMBER),
        ("Dataset parse failure",  "pandas.ParserError",        "Skip dataset, log warning, continue gather()",  TEAL),
        ("Gene mapping miss",      "KeyError / empty DataFrame","Return empty result for gene, mark as 'no map'", BLUE),
        ("LLM API error",          "mistralai.APIError",        "Return error event via SSE, user sees message",  CORAL),
        ("Cox PH non-convergence", "ConvergenceWarning",        "Skip gene, log with dataset + gene ID",          PURPLE),
        ("JWT decode failure",     "jwt.InvalidTokenError",     "Raise HTTPException(401) — never swallow auth",  CORAL),
    ]
    rh = 0.68
    rg = 0.08
    for i, (err, exc, action, col) in enumerate(errors):
        ry = 1.48 + i * (rh + rg)
        _rect(slide, rx, ry, rw, rh, CARD, f"ehR{i}")
        _rect(slide, rx, ry, 0.07, rh, col, f"ehRA{i}")
        tb_e = _tb(slide, rx + 0.14, ry + 0.05, 1.80, 0.28)
        _para(tb_e.text_frame, err, 10, col, bold=True)
        tb_x = _tb(slide, rx + 0.14, ry + 0.36, 1.80, 0.26)
        _para(tb_x.text_frame, exc, 9, MUTED)
        tb_a = _tb(slide, rx + 2.05, ry + 0.10, rw - 2.15, 0.50)
        _para(tb_a.text_frame, action, 10, LIGHT)

    _rect(slide, rx, 5.88, rw, 0.44, TEAL_DARK, "ehNote")
    tb_n = _tb(slide, rx + 0.12, 5.94, rw - 0.20, 0.30)
    _para(tb_n.text_frame, "grep WARNING geo_logs/app.log — see all recoverable errors by type and frequency", 10, TEAL)


# ── Main builder ───────────────────────────────────────────────────────────────

def create_presentation(output_path: Path, screenshots_dir: Path) -> None:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SECTION 1: Hook & Problem ──────────────────────────────────────────────
    # Slide 1
    add_title_slide(
        prs,
        "GEO Survival Analysis",
        "Cross-Dataset Cancer Biomarker Discovery from NCBI GEO",
        "PyData Yerevan · Software Engineering Section  ·  FastAPI · lifelines · pydantic-ai · React",
    )
    # Slide 2
    add_combined_intro_slide(prs)

    # ── SECTION 2: Domain Primer ───────────────────────────────────────────────
    # Slide 3
    add_survival_primer_slide(prs)
    # Slide 4
    add_km_curve_slide(prs)
    # Slide 5
    add_hr_slide(prs)
    # Slide 6
    add_meta_analysis_slide(prs)

    # ── SECTION 3: What We Built ───────────────────────────────────────────────
    # Slide 7
    add_comparison_table_slide(prs)

    # ── SECTION 3b: Clinical Reality (predictive value + responsible use) ──────
    # Slide 7a: predictive biomarkers + benefit for real patients
    add_two_col_slide(
        prs,
        "Predictive Biomarkers: Benefit for Real Patients",
        left_header="FROM PROGNOSTIC TO PREDICTIVE",
        left_header_color=TEAL,
        left_bullets=[
            "Prognostic = who is high-risk. Predictive = whose outcome changes with treatment",
            "Tested with an  expression × treatment  interaction in a Cox model — per-arm HRs + interaction p, aggregated across independent GEO cohorts",
            "Validated multi-gene signature → high / intermediate / low-risk groups (Harrell's C-index, cross-cohort) — stratified-medicine logic à la Oncotype DX",
            "Single-sample scoring → one tumour profile → reference risk group + advisory, evidence-grounded treatments to discuss",
        ],
        right_header="WHY A PATIENT BENEFITS",
        right_bullets=[
            "Risk stratification can spare low-risk patients overtreatment and flag high-risk for escalation — discussed with the care team",
            "Treatment suggestions grounded in CIViC / DGIdb biomarker→therapy evidence + outcomes from GEO cohorts where the treatment was documented",
            "A gene significant in 8 independent cohorts carries far more weight than one significant in a single study",
        ],
        accent=TEAL,
        bottom_badge="Advisory & hypothesis-generating — suggestions to discuss with the tumour board, validated prospectively",
    )
    # Slide 7b: what "predictive" means — and what it doesn't
    add_two_col_slide(
        prs,
        "What ‘Predictive’ Means Here — and What It Doesn't",
        left_header="WHAT WE CLAIM — HONESTLY",
        left_header_color=TEAL,
        left_bullets=[
            "Surfaces treatment-effect-modifying expression biomarkers validated across cohorts",
            "Turns a tumour profile into advisory treatment options to consider",
            "Every output carries a research-use, validate-prospectively label",
        ],
        right_header="WHAT WE NEVER CLAIM",
        right_bullets=[
            "Not a validated companion diagnostic — never “this patient will respond to drug X”",
            "Not a prescription — advisory “to discuss” only; directive/prescribing language is stripped at the API layer",
            "Not de-novo drug-response modelling — only documented associations",
            "Not a clinical decision-making device — clear of FDA-CDS / EU-MDR territory",
        ],
        accent=CORAL,
        bottom_badge="Research use only · hypothesis-generating · not a prescription, not a guarantee of response, not a clinical decision device",
    )
    # Slide 7c: limitations in real clinical practice
    add_gap_cards_slide(
        prs,
        "Limitations in Real Clinical Practice",
        cards=[
            ("HIPAA / PHI",
             "Patient expression submitted for scoring is held in-memory only — never persisted, never logged; no PII in logs",
             PURPLE),
            ("Research use only",
             "Not FDA-cleared, not a CDS device; complements companion-diagnostic tools, does not replace them",
             CORAL),
            ("Underpowered interactions",
             "Interaction tests are weak in small per-arm cohorts; predictive claims must state this",
             AMBER),
            ("Non-standardized annotations",
             "GEO treatment metadata is inconsistent; cohort comparability is limited",
             BLUE),
            ("Cross-platform normalization",
             "Rank / z-within-cohort normalization is mandatory — a documented method limitation",
             TEAL),
            ("Validate prospectively",
             "All findings are hypothesis-generating; clinical use requires independent prospective validation",
             PURPLE),
        ],
        accent=CORAL,
    )

    # ── SECTION 4: Backend Deep Dive ──────────────────────────────────────────
    add_backend_arch_visual_slide(prs)  # Slide 8: visual routes → services schema
    add_fastapi_slide(prs)              # Slide 9
    add_pydantic_slide(prs)             # Slide 10
    add_asyncio_slide(prs)              # Slide 11
    add_pydantic_ai_slide(prs)          # Slide 12
    add_lifelines_slide(prs)            # Slide 13
    add_sqlalchemy_slide(prs)           # Slide 14
    add_mistral_uv_slide(prs)           # Slide 15

    # ── SECTION 5: Frontend Deep Dive ─────────────────────────────────────────
    add_frontend_arch_visual_slide(prs)  # Slide 16: visual component tree + Redux + API
    add_react_slide(prs)                 # Slide 17
    add_redux_slide(prs)                 # Slide 18
    add_sse_frontend_slide(prs)          # Slide 19
    add_recharts_slide(prs)              # Slide 20

    # ── SECTION 6: Pipeline & Services ────────────────────────────────────────
    # Slide 21: 6-step workflow overview
    add_workflow_slide(
        prs,
        "The 6-Step Workflow We Automate",
        steps=[
            ("1", "Search Datasets",
             "Query NCBI Entrez eSearch with disease + tissue keywords — returns GSE accession IDs"),
            ("2", "Parse Formats",
             "Download SOFT/miniML files; parse tabular blocks; extract expression matrix + sample metadata"),
            ("3", "Map Probe → Gene",
             "Fetch GPL annotation; detect probe_id → gene_symbol column; aggregate multi-probe; cache as Parquet"),
            ("4", "Find Survival Data",
             "Detect OS time + event columns via regex; pydantic-ai handles edge cases; align to sample IDs"),
            ("5", "Run Cox + KM",
             "Median-split expression per gene → high/low groups; CoxPHFitter; log-rank; HR + 95% CI per dataset"),
            ("6", "Synthesise Across Studies",
             "Rank genes by # significant cohorts; pool HRs; report I² heterogeneity; export CSV + Methods text"),
        ],
        bottom_text="GEO Survival Analysis automates all 6 steps  ·  under 5 minutes  ·  no coding required",
    )
    add_geoclient_slide(prs)          # Slide 22
    add_geoloader_slide(prs)          # Slide 23
    add_orchestrator_slide(prs)       # Slide 24
    add_ranking_service_slide(prs)    # Slide 25
    add_gene_babel_slide(prs)         # Slide 26
    add_gene_cache_slide(prs)         # Slide 27

    # Slide 28: Statistical engine overview
    add_two_col_slide(
        prs,
        "Statistical Engine: lifelines",
        left_header="METHODS",
        left_header_color=PURPLE,
        left_bullets=[
            "Median split — expression → high / low groups per dataset",
            "Log-rank test — non-parametric group comparison (p-value)",
            "Cox PH regression — hazard ratio + 95% confidence interval",
            "KaplanMeierFitter — time-to-event survival probability curves",
            "Cross-cohort ranking — genes sorted by n_significant, then avg p",
            "Heterogeneity — Cochran Q statistic + I² on forest plots",
        ],
        right_header="DATA FLOW",
        right_bullets=[
            "INPUT: GSE series_matrix.txt.gz",
            "→ probe_id × sample expression matrix  (pandas DataFrame)",
            "→ probe → gene symbol  (GPL annotation, parquet cache)",
            "→ high / low groups  (median split per dataset)",
            "→ Cox PH fit  →  HR, p-value, 95% CI",
            "→ KM fitter  →  survival probability curves",
            "OUTPUT: ranked gene list  (cross-dataset consistency)",
        ],
        accent=PURPLE,
        bottom_badge="Significance: p < 0.05  ·  min_occurrence: 2 datasets  ·  ~600 COSMIC cancer genes filter",
    )
    add_survival_service_slide(prs)   # Slide 29

    # ── SECTION 7: AI Chat ─────────────────────────────────────────────────────
    # Slide 30: AI Chat overview
    add_architecture_slide(
        prs,
        "AI Chat: pydantic-ai + RAG + Domain Score",
        layers=[
            ("User",
             ["React Chat UI", "Natural language question", "Streaming response (SSE)"],
             BLUE),
            ("Agent",
             ["pydantic-ai  Agent", "Mistral (default)  ·  Claude Haiku (optional)",
              "Dynamic system prompt", "History trim  ≤ 16K chars"],
             TEAL),
            ("5 Tools",
             ["search_known_datasets", "search_geo_datasets",
              "get_gene_info", "estimate_query", "get_user_recent_results"],
             AMBER),
            ("RAG + Data",
             ["GEO dataset cache (parquet)", "NCBI Entrez  (live)",
              "Mistral embed 1024-dim  (numpy cosine)", "SQLite results history"],
             PURPLE),
        ],
        accent=TEAL,
    )
    add_chat_tools_slide(prs)         # Slide 31
    add_rag_slide(prs)                # Slide 32
    add_domain_score_slide(prs)       # Slide 33

    # ── SECTION 8: Cross-Cutting Concerns ─────────────────────────────────────
    add_sse_pipeline_slide(prs)       # Slide 34
    add_jwt_slide(prs)                # Slide 35
    add_data_pipeline_slide(prs)      # Slide 36
    add_caching_layers_slide(prs)     # Slide 37
    add_asyncio_gather_slide(prs)     # Slide 38
    add_error_handling_slide(prs)     # Slide 39

    # ── SECTION 9: Engineering Decisions ──────────────────────────────────────
    add_db_comparison_slide(prs)      # Slide 40

    # Slide 41: Deployment
    add_two_col_slide(
        prs,
        "Deployment Pipeline",
        left_header="MULTI-STAGE DOCKER BUILD",
        left_header_color=TEAL,
        left_bullets=[
            "Stage 1: node:22-alpine — npm ci + vite build  (React SPA → /dist)",
            "Stage 2: uv:python3.13 — uv sync --frozen --no-dev  (.venv)",
            "Stage 3: python:3.13-slim runtime — copy .venv + frontend/dist",
            "CMD: alembic upgrade head && uvicorn app.main:app --port 8000",
            "GitHub Actions: push → build → push ghcr.io → SSH deploy",
        ],
        right_header="HETZNER CX32  ·  CADDY  ·  DOCKER COMPOSE",
        right_bullets=[
            "Caddy: automatic TLS + flush_interval=-1  (enables SSE streaming)",
            "SQLite default; DATABASE_URL env var → PostgreSQL (no code change)",
            "Named volumes: app_data · platform_mappings · geo_logs",
            "Why Caddy: zero-config TLS + native SSE, one-line Caddyfile",
            "Cost: ~$10/mo  (CX32 — 4 vCPU, 8 GB RAM, 80 GB NVMe)",
        ],
        accent=TEAL,
        bottom_badge="git push → live in ~3 min  ·  zero-downtime docker compose up -d",
    )

    # ── SECTION 10: Demo ───────────────────────────────────────────────────────
    # Slide 42
    add_image_slide(
        prs,
        "Demo: Natural Language Query",
        _find_screenshot(screenshots_dir, "starting_page.png"),
        'Type a natural language query — e.g. "What genes predict poor survival in triple-negative breast cancer?"',
    )
    # Slide 43
    add_image_slide(
        prs,
        "Demo: Volcano Plot & Gene Results",
        _find_screenshot(screenshots_dir, "volcano_plot.png"),
        "Ranked genes with hazard ratios, p-values, and risk direction consistency across datasets",
    )
    # Slide 44
    add_image_slide(
        prs,
        "Demo: Kaplan-Meier Survival Curves",
        _find_screenshot(screenshots_dir, "kaplan_meier_curves.png", "km_curves.png"),
        "Compare survival probability between high and low expression groups — interpretable by clinicians",
    )

    # ── SECTION 11: Lessons Learned ───────────────────────────────────────────
    # Slide 45
    add_gap_cards_slide(
        prs,
        "What We Learned",
        cards=[
            ("asyncio.gather",
             "Parallel downloads cut wall time from 20 min → 2 min — the single biggest engineering win",
             TEAL),
            ("Hybrid metadata detection",
             "Regex catches ~70% of survival columns; pydantic-ai structured output handles the rest",
             PURPLE),
            ("Disk-first caching",
             "Parquet in platform_mappings/ survives restarts; OrderedDict LRU is a fast in-memory layer on top",
             BLUE),
            ("uv in Docker",
             "uv sync --frozen: reproducible, layer-cached, no requirements.txt drift — 10× faster builds",
             AMBER),
            ("RAG at query time",
             "Embedding + cosine similarity retrieves live context; agent cites real data, no hallucinations",
             TEAL),
            ("pydantic-ai is lightweight",
             "No LangChain overhead — structured output, tool calling, streaming all work natively",
             CORAL),
        ],
        accent=TEAL,
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved {output_path}  ({output_path.stat().st_size // 1024} KB,  48 slides)")


if __name__ == "__main__":
    project_root    = Path(__file__).parent.parent
    screenshots_dir = project_root / "presentations" / "screenshots"
    output_path     = project_root / "presentations" / "app_presentation.pptx"

    if not screenshots_dir.exists():
        print(f"Error: Screenshots directory not found: {screenshots_dir}")
        raise SystemExit(1)

    create_presentation(output_path, screenshots_dir)
